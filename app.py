import streamlit as st
import pandas as pd
import plotly.express as px
import plotly.graph_objects as go
import os
import urllib.request
import json
import io
from streamlit_oauth import OAuth2Component
from ops_dashboard import show_dashboard
# ==========================================
# PAGE CONFIGURATION & CUSTOM CSS
# ==========================================
st.set_page_config(
    page_title="Impact Analytics Dashboard",
    layout="wide",
    initial_sidebar_state="expanded",
)

st.markdown("""
<style>
    [data-testid="stMetricValue"] { font-size: 2rem; font-weight: 700; }
    .stTabs [data-baseweb="tab-list"] { gap: 24px; }
    .stTabs [data-baseweb="tab"] {
        height: 50px; white-space: pre-wrap; background-color: transparent;
        border-radius: 4px 4px 0px 0px; padding-top: 10px; padding-bottom: 10px;
    }
    .student-card {
        background: #f8f9fa; border-radius: 10px; padding: 16px 20px;
        border-left: 5px solid #0094c9; margin-bottom: 12px;
    }
</style>
""", unsafe_allow_html=True)

# ==========================================
# MODULE-LEVEL CONSTANTS
# ==========================================
COLOR_MAP    = {'Baseline': '#636EFA', 'Endline': '#00CC96'}
AY_COLOR_MAP = {'AY24-25': '#636EFA', 'AY25-26': '#00CC96'}
RISE_COLORS  = {
    "Reviving":   "#f27c48",
    "Initiating": "#0094c9",
    "Shaping":    "#00964d",
    "Evolving":   "#ed1c2d",
}
RISE_COLORS_LONG = RISE_COLORS
RISE_ORDER = ["Reviving", "Initiating", "Shaping", "Evolving"]
AY_ORDER   = ["AY24-25", "AY25-26"]

# AY 25-26 uses 2-letter state abbreviations; normalise to full names
STATE_ABBR = {
    "AP": "Andhra Pradesh", "AS": "Assam",         "BR": "Bihar",
    "GJ": "Gujarat",        "JH": "Jharkhand",     "KA": "Karnataka",
    "MH": "Maharashtra",    "MP": "Madhya Pradesh", "OD": "Odisha",
    "RJ": "Rajasthan",      "TN": "Tamil Nadu",     "TS": "Telangana",
    "UK": "Uttarakhand",    "UP": "Uttar Pradesh",
}

# Subject name normalisation across both workbooks:
#   "Maths"                  → "Math"              (AY24-25 spelling)
#   "Basic Digital Literacy" → "Digital Literacy"  (AY24-25 long form)
#   "DL"                     → "Digital Literacy"  (AY25-26 abbreviation)
SUBJECT_NORM = {
    "Maths":                  "Math",
    "Basic Digital Literacy": "Digital Literacy",
    "DL":                     "Digital Literacy",
}

# Sentinel used in clean_sheet for rows that have no Gender column (AY 24-25)
NO_GENDER = "__NO_GENDER__"


# ==========================================
# AUTHENTICATION GATEKEEPER
# ==========================================
try:
    CLIENT_ID     = st.secrets["GOOGLE_CLIENT_ID"]
    CLIENT_SECRET = st.secrets["GOOGLE_CLIENT_SECRET"]
except (FileNotFoundError, KeyError):
    st.error(
        "Missing `.streamlit/secrets.toml` or Streamlit Cloud Secrets. "
        "Please ensure your Google Client ID and Secret are configured."
    )
    st.stop()

AUTHORIZE_URL    = "https://accounts.google.com/o/oauth2/v2/auth"
TOKEN_URL        = "https://oauth2.googleapis.com/token"
REVOKE_TOKEN_URL = "https://oauth2.googleapis.com/revoke"

for _k, _v in [("logged_in_email", None), ("user_first_name", "User")]:
    if _k not in st.session_state:
        st.session_state[_k] = _v

if not st.session_state["logged_in_email"]:
    _, col2, _ = st.columns(3)
    with col2:
        st.write(""); st.write("")
        try:
            st.image("evidyaloka_logo.png", width=320)
        except Exception:
            pass
        st.markdown(
            "<h2 style='text-align:center;color:#0094c9;'>Student Analytics Portal</h2>",
            unsafe_allow_html=True,
        )
        st.markdown(
            "<p style='text-align:center;'>Sign in with your @evidyaloka.org email.</p>",
            unsafe_allow_html=True,
        )
        st.markdown("---")
        oauth2 = OAuth2Component(CLIENT_ID, CLIENT_SECRET, AUTHORIZE_URL, TOKEN_URL, REVOKE_TOKEN_URL)
        result = oauth2.authorize_button(
            name="Sign in with Google",
            icon="https://upload.wikimedia.org/wikipedia/commons/5/53/Google_%22G%22_Logo.svg",
            redirect_uri="https://ev-assessments.streamlit.app",
            scope="openid email profile",
            key="google_login",
            use_container_width=True,
        )
        if result and "token" in result:
            id_token = result["token"]["id_token"]
            if isinstance(id_token, list):
                id_token = id_token[0]
            verify_url = f"https://oauth2.googleapis.com/tokeninfo?id_token={id_token}"
            try:
                with urllib.request.urlopen(verify_url) as resp:
                    user_info = json.loads(resp.read().decode())
                st.session_state["logged_in_email"]  = user_info.get("email")
                st.session_state["user_first_name"] = user_info.get("given_name", "User")
                st.rerun()
            except Exception as e:
                st.error(f"Error verifying login with Google: {e}")
                st.stop()
    st.stop()


# ==========================================
# SHARED SIDEBAR FILTER BUILDER
# ==========================================
def build_filter_sidebar(df: pd.DataFrame, key_prefix: str):
    """
    Cascading State → Donor → Centre → Subject → Grade → Gender filters.

    BUG FIX: The gender filter now only filters rows that HAVE a known gender
    (i.e. rows from AY 25-26). Rows where Gender == NO_GENDER (all AY 24-25
    rows, which have no gender column) are always kept regardless of the
    gender selection. This prevents the gender multiselect from silently
    wiping out all AY 24-25 data.
    """
    sel = {}
    st.sidebar.header("🎯 Global Filters")

    def _opts(series):
        return ["All"] + sorted(series.dropna().astype(str).unique())

    # State
    states = _opts(df["State"]) if "State" in df.columns else ["All"]
    s = st.sidebar.selectbox("Select State", states, key=f"{key_prefix}_state")
    sel["state"] = s
    dff = df[df["State"].astype(str) == s].copy() if s != "All" else df.copy()

    # Donor
    donors = _opts(dff["Donor"]) if "Donor" in dff.columns else ["All"]
    d = st.sidebar.selectbox("Select Donor", donors, key=f"{key_prefix}_donor")
    sel["donor"] = d
    if d != "All":
        dff = dff[dff["Donor"].astype(str) == d]

    # Centre
    centres = _opts(dff["Centre Name"]) if "Centre Name" in dff.columns else ["All"]
    c = st.sidebar.selectbox("Select Centre", centres, key=f"{key_prefix}_centre")
    sel["centre"] = c
    if c != "All":
        dff = dff[dff["Centre Name"].astype(str) == c]

    # Subject
    subjects = _opts(dff["Subject"]) if "Subject" in dff.columns else ["All"]
    sub = st.sidebar.selectbox("Select Subject", subjects, key=f"{key_prefix}_subject")
    sel["subject"] = sub
    if sub != "All":
        dff = dff[dff["Subject"].astype(str) == sub]

    # Grade
    if key_prefix == "long":
        # Longitudinal: derive cohort from AY24-25 grades, then track into AY25-26
        base_yr = dff[dff["Academic Year"] == "AY24-25"] if "Academic Year" in dff.columns else dff
        grades = sorted(base_yr["Grade"].dropna().astype(str).unique()) if "Grade" in base_yr.columns else []
        sel_grades = st.sidebar.multiselect(
            "Select AY 24-25 Grade (Cohort)", options=grades, default=grades, key=f"{key_prefix}_grade",
            help="Tracks students from their AY 24-25 grade into their promoted grade in AY 25-26.",
        )
        if sel_grades:
            cohort_ids = base_yr[base_yr["Grade"].astype(str).isin(sel_grades)]["Student ID"].unique()
            dff = dff[dff["Student ID"].isin(cohort_ids)]
        else:
            dff = dff.iloc[0:0]
    else:
        grades = sorted(dff["Grade"].dropna().astype(str).unique()) if "Grade" in dff.columns else []
        sel_grades = st.sidebar.multiselect(
            "Select Grade(s)", options=grades, default=grades, key=f"{key_prefix}_grade"
        )
        dff = dff[dff["Grade"].astype(str).isin(sel_grades)] if sel_grades else dff.iloc[0:0]
    sel["grades"] = sel_grades

    # ── Gender filter ────────────────────────────────────────────────────────
    # CRITICAL FIX: AY 24-25 rows carry Gender == NO_GENDER (a sentinel, not a
    # real gender value). We must NOT drop those rows when the user selects
    # genders — the filter should only restrict rows that have a real gender.
    # Implementation: keep (Gender in selected) OR (Gender == NO_GENDER).
    if "Gender" in dff.columns:
        real_gender_rows = dff[~dff["Gender"].astype(str).isin([NO_GENDER, "nan", "none", "null", ""])]
        genders = sorted(real_gender_rows["Gender"].astype(str).unique())
        if genders:
            sel_g = st.sidebar.multiselect(
                "Select Gender(s) — applies to AY 25-26 only",
                options=genders, default=genders, key=f"{key_prefix}_gender",
            )
            # Keep rows that match selection OR have no gender data (AY 24-25)
            mask = (
                dff["Gender"].astype(str).isin(sel_g) |
                dff["Gender"].astype(str).isin([NO_GENDER, "nan", "none", "null", ""])
            )
            dff = dff[mask]
            sel["genders"] = sel_g
        else:
            sel["genders"] = []
    else:
        sel["genders"] = []

    return dff, sel


# ==========================================
# APP ROUTER / HOME PAGE
# ==========================================
if "current_page" not in st.session_state:
    st.session_state["current_page"] = "home"

if st.session_state["current_page"] == "home":
    st.write(""); st.write("")
    st.title(f"👋 Welcome, {st.session_state['user_first_name']}!")
    st.markdown(
        "<p style='color:gray;font-size:1.1em;'>Select an application below to continue.</p>",
        unsafe_allow_html=True,
    )
    st.markdown("---")
    col1, col2, col3, col4, *_ = st.columns(4)
    with col1:
        st.markdown("<h1 style='text-align:center;font-size:4rem;'>📈</h1>", unsafe_allow_html=True)
        if st.button("Impact Analytics Dashboard", use_container_width=True):
            st.session_state["current_page"] = "dashboard"; st.rerun()
    with col2:
        st.markdown("<h1 style='text-align:center;font-size:4rem;'>🏛️</h1>", unsafe_allow_html=True)
        if st.button("Longitudinal Analysis", use_container_width=True):
            st.session_state["current_page"] = "longitudinal"; st.rerun()
    with col3:
        st.markdown("<h1 style='text-align:center;font-size:4rem;'>🏢</h1>", unsafe_allow_html=True)
        if st.button("Operations & Impact Command Center", use_container_width=True):
            st.session_state["current_page"] = "ops"; st.rerun()
    with col4:
        st.markdown("<h1 style='text-align:center;font-size:4rem;'>📚</h1>", unsafe_allow_html=True)
        if st.button("Topic & SubTopic Analytics", use_container_width=True):
            st.session_state["current_page"] = "topics"; st.rerun()
    st.stop()


# ==========================================
# OPS DASHBOARD MODULE
# ==========================================
if st.session_state["current_page"] == "ops":
    # Sidebar nav (mirrors the other modules)
    with st.sidebar:
        try:
            st.image("evidyaloka_logo.png", width=273)
        except Exception:
            pass
        st.success(f"👤 **{st.session_state['user_first_name']}**")
        nc1, nc2 = st.columns(2)
        with nc1:
            if st.button("🏠 Home", use_container_width=True, key="nav_home_ops"):
                st.session_state["current_page"] = "home"; st.rerun()
        with nc2:
            if st.button("Sign Out", use_container_width=True, key="signout_ops"):
                st.session_state.update({
                    "logged_in_email": None,
                    "user_first_name": "User",
                    "current_page": "home",
                })
                st.rerun()

    # Import and render — ops_dashboard.py must sit alongside app.py
    try:
        from ops_dashboard import render_ops_dashboard
        render_ops_dashboard()
    except ImportError:
        st.error(
            "⚠️ `ops_dashboard.py` not found. "
            "Place it in the same folder as `app.py` and restart the app."
        )
    st.stop()


# ==========================================
# TOPIC & SUBTOPIC ANALYTICS MODULE
# ==========================================
if st.session_state["current_page"] == "topics":
    # Sidebar nav (mirrors the other modules)
    with st.sidebar:
        try:
            st.image("evidyaloka_logo.png", width=273)
        except Exception:
            pass
        st.success(f"👤 **{st.session_state['user_first_name']}**")
        nc1, nc2 = st.columns(2)
        with nc1:
            if st.button("🏠 Home", use_container_width=True, key="nav_home_topics"):
                st.session_state["current_page"] = "home"; st.rerun()
        with nc2:
            if st.button("Sign Out", use_container_width=True, key="signout_topics"):
                st.session_state.update({
                    "logged_in_email": None,
                    "user_first_name": "User",
                    "current_page": "home",
                })
                st.rerun()

    # Import and render — topic_dashboard.py must sit alongside app.py
    try:
        from topic_dashboard import render_topic_dashboard
        render_topic_dashboard()
    except ImportError:
        st.error(
            "⚠️ `topic_dashboard.py` not found. "
            "Place it in the same folder as `app.py` and restart the app."
        )
    st.stop()


# ==========================================
# DATA LOADING — LONGITUDINAL MODULE
# ==========================================
FILE_24 = "EL-BL-Data-AY-24-25.xlsx"
FILE_25 = "AltenUP-WB-BL-EL-Data-AY25-26.xlsx"


@st.cache_data(show_spinner=False)
def load_multi_year_data(src_24, src_25):
    """
    Loads both workbooks and returns a single normalised DataFrame.

    Key normalisations:
    - State abbreviations → full names (AY25-26 uses 2-letter codes)
    - "Maths" → "Math" (AY24-25 spelling)
    - Gender: title-cased and stripped; rows WITHOUT a gender column
      (all AY24-25 rows) receive Gender = NO_GENDER sentinel so the
      sidebar gender filter can distinguish them from genuinely unknown values.
    - Student ID → clean string via Int64 path (numeric) or upper-stripped (alphanumeric)
    """
    def clean_sheet(df_raw, year: str, period: str) -> pd.DataFrame:
        if df_raw.empty:
            return pd.DataFrame()
        df = df_raw.copy()

        # Rubrics → Category
        if "Rubrics" in df.columns and "Category" not in df.columns:
            df.rename(columns={"Rubrics": "Category"}, inplace=True)

        # Robust Student ID: numeric → Int64 → str; alphanumeric → upper + strip
        raw_ids     = df.get("Student ID", pd.Series(dtype=object)).copy()
        numeric_ids = pd.to_numeric(raw_ids, errors="coerce")
        is_num      = numeric_ids.notna()
        df["Student ID"] = pd.NA
        df.loc[is_num,  "Student ID"] = numeric_ids[is_num].astype("Int64").astype(str).str.strip()
        df.loc[~is_num, "Student ID"] = raw_ids[~is_num].astype(str).str.strip().str.upper()
        df.loc[df["Student ID"].astype(str).str.upper().isin(
            ["NA", "NAN", "NONE", "NULL", "<NA>", ""]
        ), "Student ID"] = pd.NA

        # Normalise State (expand 2-letter codes used by AY25-26)
        if "State" in df.columns:
            df["State"] = df["State"].astype(str).str.strip().map(lambda x: STATE_ABBR.get(x, x))

        # Normalise Subject
        if "Subject" in df.columns:
            df["Subject"] = df["Subject"].astype(str).str.strip().map(lambda x: SUBJECT_NORM.get(x, x))

        # Gender handling — CRITICAL:
        # AY24-25 workbook has NO Gender column. We use the NO_GENDER sentinel
        # so the sidebar gender filter can keep these rows unconditionally.
        if "Gender" in df.columns:
            df["Gender"] = (
                df["Gender"].astype(str).str.strip().str.title()
                .str.replace(r"\s+", "", regex=True)   # "Boy " → "Boy"
            )
            # Replace residual junk values with sentinel
            df.loc[df["Gender"].str.lower().isin(["nan", "none", "null", ""]), "Gender"] = NO_GENDER
        else:
            df["Gender"] = NO_GENDER

        # Other string columns
        for col in ["Centre Name", "Donor", "Category"]:
            if col in df.columns:
                df[col] = df[col].astype(str).str.strip()

        if "Grade" in df.columns:
            df["Grade"] = df["Grade"].astype(str).str.replace(r"\.0$", "", regex=True)

        df["Obtained Marks"] = pd.to_numeric(df.get("Obtained Marks"), errors="coerce")
        df["Total Marks"]    = pd.to_numeric(df.get("Total Marks"),    errors="coerce")
        df["Pct Score"]      = (df["Obtained Marks"] / df["Total Marks"] * 100).round(1)

        df["Academic Year"] = year
        df["Period"]        = period
        df["Timepoint"]     = f"{year} {period}"

        keep = [
            "State", "Centre Name", "Donor", "Subject", "Grade", "Student ID",
            "Gender", "Obtained Marks", "Total Marks", "Pct Score",
            "Category", "Academic Year", "Period", "Timepoint",
        ]
        return df[[c for c in keep if c in df.columns]].dropna(subset=["Student ID", "Obtained Marks"])

    try:
        xls24   = pd.ExcelFile(src_24)
        xls25   = pd.ExcelFile(src_25)
        sheets24 = xls24.sheet_names
        sheets25 = xls25.sheet_names

        bl24 = clean_sheet(pd.read_excel(src_24, sheet_name=sheets24[0]),                              "AY24-25", "Baseline")
        el24 = clean_sheet(pd.read_excel(src_24, sheet_name=sheets24[1] if len(sheets24) > 1 else 0), "AY24-25", "Endline")
        bl25 = clean_sheet(pd.read_excel(src_25, sheet_name=sheets25[0]),                              "AY25-26", "Baseline")
        el25 = clean_sheet(pd.read_excel(src_25, sheet_name=sheets25[1] if len(sheets25) > 1 else 0), "AY25-26", "Endline")

        combined = pd.concat([bl24, el24, bl25, el25], ignore_index=True)
        return combined
    except Exception as e:
        st.error(f"Error loading multi-year data: {e}")
        return None


# ==========================================
# LONGITUDINAL ANALYSIS MODULE
# ==========================================
if st.session_state["current_page"] == "longitudinal":
    st.title("🏛️ Strategic Longitudinal Analysis")
    st.markdown(
        "<p style='color:gray;font-size:1.1em;'>Year-over-Year Trajectories & Student Deep-Dives "
        "(AY 24-25 vs AY 25-26 — Endline to Endline)</p>",
        unsafe_allow_html=True,
    )
    st.markdown("---")

    # Sidebar nav
    with st.sidebar:
        try:
            st.image("evidyaloka_logo.png", width=273)
        except Exception:
            pass
        st.success(f"👤 **{st.session_state['user_first_name']}**")
        nc1, nc2 = st.columns(2)
        with nc1:
            if st.button("🏠 Home", use_container_width=True, key="nav_home_long"):
                st.session_state["current_page"] = "home"; st.rerun()
        with nc2:
            if st.button("Sign Out", use_container_width=True, key="signout_long"):
                st.session_state.update({"logged_in_email": None, "user_first_name": "User", "current_page": "home"})
                st.rerun()
        st.markdown("---")
        st.info(
            "💡 **YoY logic:** All comparison charts use **Endline scores only** "
            "(AY24-25 Endline vs AY25-26 Endline). "
            "Tabs 1–3 additionally require matched Student IDs across both years."
        )

    # File resolution
    src_24, src_25 = FILE_24, FILE_25
    if not (os.path.exists(FILE_24) and os.path.exists(FILE_25)):
        st.warning("Data files not found on disk. Upload them below.")
        up24 = st.file_uploader(f"Upload {FILE_24}", type=["xlsx"], key="up_24")
        up25 = st.file_uploader(f"Upload {FILE_25}", type=["xlsx"], key="up_25")
        if not (up24 and up25):
            st.info("Please upload both files to continue.")
            st.stop()
        src_24, src_25 = up24, up25

    with st.spinner("Loading and harmonising multi-year data…"):
        df_long = load_multi_year_data(src_24, src_25)

    if df_long is None or df_long.empty:
        st.error("Could not load data. Check the uploaded files.")
        st.stop()

    # Apply sidebar filters
    with st.sidebar:
        filtered_df_long, long_sel = build_filter_sidebar(df_long, key_prefix="long")

    if filtered_df_long.empty:
        st.warning("⚠️ No data for the selected filters. Please adjust your criteria.")
        st.stop()

    # ── Core slices ─────────────────────────────────────────────────────────
    # BUG FIX: Do NOT assign pd.Categorical to the "Academic Year" column of
    # df_endlines at this point. Apply it only immediately before plotting,
    # and cast back to str for the actual chart x-axis argument so Plotly
    # receives plain strings (preserving correct ordering via category_orders).
    df_el24     = filtered_df_long[filtered_df_long["Timepoint"] == "AY24-25 Endline"].copy()
    df_el25     = filtered_df_long[filtered_df_long["Timepoint"] == "AY25-26 Endline"].copy()
    df_endlines = filtered_df_long[filtered_df_long["Period"] == "Endline"].copy()
    # Keep Academic Year as plain string — ordering is enforced via
    # category_orders= in every px.line / px.bar call.

    # Retained cohort (Student IDs present in both endlines)
    ids_24     = set(df_el24["Student ID"].dropna().unique())
    ids_25     = set(df_el25["Student ID"].dropna().unique())
    retained   = ids_24 & ids_25
    no_overlap = len(retained) == 0
    df_ret_24  = df_el24[df_el24["Student ID"].isin(retained)] if not no_overlap else pd.DataFrame()
    df_ret_25  = df_el25[df_el25["Student ID"].isin(retained)] if not no_overlap else pd.DataFrame()

    # Debug expander
    with st.expander("🛠️ Debug: ID Matching", expanded=False):
        dc1, dc2 = st.columns(2)
        with dc1:
            st.markdown("**AY 24-25 Endline**")
            st.write(f"Rows: `{len(df_el24)}` | Unique IDs: `{len(ids_24)}`")
            st.dataframe(
                df_el24["Student ID"].dropna().drop_duplicates().head(5)
                .reset_index(drop=True).rename("Student ID"),
                use_container_width=True,
            )
        with dc2:
            st.markdown("**AY 25-26 Endline**")
            st.write(f"Rows: `{len(df_el25)}` | Unique IDs: `{len(ids_25)}`")
            st.dataframe(
                df_el25["Student ID"].dropna().drop_duplicates().head(5)
                .reset_index(drop=True).rename("Student ID"),
                use_container_width=True,
            )
        st.write(f"**Intersection size:** `{len(retained)}`")
        if no_overlap and ids_24 and ids_25:
            one24, one25 = next(iter(ids_24)), next(iter(ids_25))
            st.write(f"AY24 sample → `{repr(one24)}` (len={len(one24)})")
            st.write(f"AY25 sample → `{repr(one25)}` (len={len(one25)})")

    # ── TABS ────────────────────────────────────────────────────────────────
    (mig_tab, sub_tab, gen_tab,
     subj_tab, geo_tab, centre_tab, student_tab) = st.tabs([
        "📊 Migration",
        "📚 Subject Efficacy",
        "🚻 Gender Equity",
        "📚 Subject Wise",
        "🗺️ Geographical",
        "🏫 Centre Deep Dive",
        "🔍 Student Lookup",
    ])

    # ────────────────────────────────────────────────────────────────────────
    # TAB 1: MIGRATION (retained cohort)
    # ────────────────────────────────────────────────────────────────────────
    with mig_tab:
        st.markdown("### 🧱 Structural Tier Migration (Retained Cohort Only)")
        if no_overlap:
            st.warning(
                "⚠️ No overlapping Student IDs found between AY 24-25 and AY 25-26 Endlines. "
                "Use the Subject Wise / Geographical / Centre tabs for aggregate YoY comparison, "
                "or the Student Lookup tab for individual records."
            )
        else:
            col_m1, col_m2 = st.columns([1.5, 1])
            with col_m1:
                cat24 = df_ret_24["Category"].value_counts(normalize=True).reset_index()
                cat24.columns = ["Category", "Pct"]; cat24["Year"] = "AY 24-25 Endline"
                cat25 = df_ret_25["Category"].value_counts(normalize=True).reset_index()
                cat25.columns = ["Category", "Pct"]; cat25["Year"] = "AY 25-26 Endline"
                cdf = pd.concat([cat24, cat25]); cdf["Pct"] *= 100
                fig_cat = px.bar(
                    cdf, x="Year", y="Pct", color="Category",
                    color_discrete_map=RISE_COLORS,
                    text=cdf["Pct"].apply(lambda x: f"{x:.1f}%"),
                    category_orders={
                        "Category": RISE_ORDER,
                        "Year": ["AY 24-25 Endline", "AY 25-26 Endline"],
                    },
                )
                fig_cat.update_layout(
                    barmode="stack", xaxis_title="", yaxis_title="% of Cohort",
                    plot_bgcolor="rgba(0,0,0,0)", margin=dict(l=0, r=0, t=30),
                    legend=dict(orientation="h", yanchor="top", y=-0.15, xanchor="center", x=0.5, title=""),
                )
                st.plotly_chart(fig_cat, use_container_width=True)
            with col_m2:
                st.info(
                    "Only students present in **both** endlines are shown. "
                    "Shrinking red/orange and growing green sections indicate upward migration."
                )
                try:
                    rv24 = cat24[cat24["Category"] == "Reviving"]["Pct"].values
                    rv25 = cat25[cat25["Category"] == "Reviving"]["Pct"].values
                    diff = (rv25[0] if len(rv25) else 0)*100 - (rv24[0] if len(rv24) else 0)*100
                    st.success(f"**🔍 'Reviving' proportion changed by {diff:+.1f}% YoY.**")
                    st.markdown(
                        "**💡** Excellent — the struggling base is shrinking."
                        if diff < 0 else
                        "**💡** The struggling cohort is growing. Consider targeted Tier-1 intervention."
                    )
                except Exception:
                    st.write("Insufficient category data for insights.")

    # ────────────────────────────────────────────────────────────────────────
    # TAB 2: SUBJECT EFFICACY — per-grade slopegraphs (all students, Endline only)
    # Each grade gets its own slopegraph: x-axis = AY 24-25 | AY 25-26,
    # one coloured line per subject present in that grade.
    # Uses df_endlines (not the retained cohort) so every student is included.
    # ────────────────────────────────────────────────────────────────────────
    with sub_tab:
        st.markdown("### 📈 Subject Efficacy — Per-Grade YoY Slopegraph")
        st.caption(
            "Each chart shows how subjects performed within a single grade across both academic years "
            "(AY 24-25 Endline → AY 25-26 Endline). Subjects that only appear in one year show a "
            "single dot with no connecting line — that's expected given curriculum differences."
        )

        # Build the endline-only slice with normalised subjects
        # (SUBJECT_NORM now maps Basic Digital Literacy + DL → Digital Literacy)
        df_eff = df_endlines.copy()

        # Numeric grade sort so Grade 10 doesn't precede Grade 5
        grades_eff = sorted(
            df_eff["Grade"].dropna().unique(),
            key=lambda x: int(x) if str(x).isdigit() else 99,
        )

        if df_eff.empty:
            st.info("No endline data available for the current filters.")
        else:
            # Overall insight callout (across all grades, shared subjects only)
            try:
                overall_24 = df_eff[df_eff["Academic Year"] == "AY24-25"].groupby("Subject")["Obtained Marks"].mean()
                overall_25 = df_eff[df_eff["Academic Year"] == "AY25-26"].groupby("Subject")["Obtained Marks"].mean()
                overall_delta = (overall_25 - overall_24).dropna()
                if not overall_delta.empty:
                    best_sub  = overall_delta.idxmax()
                    worst_sub = overall_delta.idxmin()
                    st.success(
                        f"**Across all grades — {best_sub}** showed the highest YoY growth "
                        f"({overall_delta[best_sub]:+.2f} pts). "
                        f"**{worst_sub}** showed the least momentum ({overall_delta[worst_sub]:+.2f} pts)."
                    )
                    st.markdown(
                        f"**💡** Cross-pollinate **{best_sub}** teaching strategies into "
                        f"**{worst_sub}** curriculum planning for the next cycle."
                    )
            except Exception:
                pass

            st.markdown("---")

            # One slopegraph per grade, laid out in pairs (2 columns)
            grade_pairs = [grades_eff[i:i+2] for i in range(0, len(grades_eff), 2)]

            for pair in grade_pairs:
                cols = st.columns(len(pair))
                for col_widget, grade in zip(cols, pair):
                    df_grade = df_eff[df_eff["Grade"] == grade]

                    # Average endline score per subject per academic year for this grade
                    grade_slope = (
                        df_grade.groupby(["Academic Year", "Subject"])["Obtained Marks"]
                        .mean().reset_index()
                    )
                    grade_slope.rename(columns={"Academic Year": "Year", "Obtained Marks": "Avg Score"}, inplace=True)
                    # Map "AY24-25" → "AY 24-25" for display consistency with the screenshot style
                    grade_slope["Year"] = grade_slope["Year"].map({"AY24-25": "AY 24-25", "AY25-26": "AY 25-26"})

                    with col_widget:
                        if grade_slope.empty:
                            st.info(f"Grade {grade}: no data.")
                        else:
                            # Count students for the subtitle
                            n_students = df_grade["Student ID"].nunique()
                            subjects_in_grade = sorted(grade_slope["Subject"].unique())

                            fig_g = px.line(
                                grade_slope,
                                x="Year", y="Avg Score", color="Subject",
                                markers=True, text="Avg Score",
                                category_orders={"Year": ["AY 24-25", "AY 25-26"]},
                                title=f"Grade {grade}  <sup>({n_students} students)</sup>",
                            )
                            fig_g.update_traces(
                                textposition="top center",
                                texttemplate="%{text:.1f}",
                                marker=dict(size=11),
                                line=dict(width=2.5),
                            )
                            fig_g.update_layout(
                                xaxis_title="",
                                yaxis_title="Avg Endline Score",
                                yaxis=dict(
                                    showgrid=True, gridcolor="lightgrey",
                                    zeroline=False, range=[0, 10.5],
                                ),
                                xaxis=dict(showgrid=False, linecolor="black"),
                                plot_bgcolor="rgba(0,0,0,0)",
                                margin=dict(l=0, r=0, t=50, b=10),
                                legend=dict(
                                    title="Subject",
                                    orientation="v",
                                    yanchor="top", y=1,
                                    xanchor="left", x=1.02,
                                    font=dict(size=11),
                                ),
                                height=380,
                            )
                            st.plotly_chart(fig_g, use_container_width=True)

                            # Compact insight: best and worst subject for this grade
                            try:
                                g24 = grade_slope[grade_slope["Year"] == "AY 24-25"].set_index("Subject")["Avg Score"]
                                g25 = grade_slope[grade_slope["Year"] == "AY 25-26"].set_index("Subject")["Avg Score"]
                                delta_g = (g25 - g24).dropna()
                                if not delta_g.empty:
                                    b = delta_g.idxmax(); w = delta_g.idxmin()
                                    st.caption(
                                        f"📈 **{b}** +{delta_g[b]:.2f} pts &nbsp;|&nbsp; "
                                        f"📉 **{w}** {delta_g[w]:+.2f} pts"
                                    )
                            except Exception:
                                pass

                st.markdown("---")

    # ────────────────────────────────────────────────────────────────────────
    # TAB 3: GENDER EQUITY (retained cohort, gender only in AY25-26)
    # ────────────────────────────────────────────────────────────────────────
    with gen_tab:
        st.markdown("### 🚻 Gender Equity — AY 25-26 Endline (Retained Cohort)")
        st.caption(
            "Gender data was not collected in AY 24-25. "
            "This tab shows the gender split for the retained cohort in their AY 25-26 Endline."
        )
        if no_overlap:
            st.warning("⚠️ No overlapping Student IDs.")
        else:
            g25_valid = df_ret_25[df_ret_25["Gender"].isin(["Boy", "Girl"])]
            if g25_valid.empty:
                st.info("No valid gender data in the retained cohort for AY 25-26.")
            else:
                col_g1, col_g2 = st.columns([1.5, 1])
                g25_avg = g25_valid.groupby("Gender")["Obtained Marks"].mean().reset_index()
                with col_g1:
                    fig_gen = px.bar(
                        g25_avg, x="Gender", y="Obtained Marks", color="Gender",
                        color_discrete_map={"Boy": "#636EFA", "Girl": "#EF553B"},
                        text=g25_avg["Obtained Marks"].apply(lambda x: f"{x:.2f}"),
                    )
                    fig_gen.update_layout(
                        xaxis_title="", yaxis_title="Avg Score (AY 25-26 Endline)",
                        plot_bgcolor="rgba(0,0,0,0)", margin=dict(l=0, r=0, t=30), showlegend=False,
                    )
                    st.plotly_chart(fig_gen, use_container_width=True)
                with col_g2:
                    try:
                        g_b = g25_avg[g25_avg["Gender"] == "Boy"]["Obtained Marks"].values[0]
                        g_g = g25_avg[g25_avg["Gender"] == "Girl"]["Obtained Marks"].values[0]
                        gap = abs(g_g - g_b)
                        leader = "Girls" if g_g > g_b else "Boys"
                        st.success(f"**🔍 Gap:** {gap:.2f} pts — **{leader}** outperform on average.")
                        st.markdown(
                            "**💡** Gap is within normal range. Monitor across future years."
                            if gap < 0.5 else
                            "**💡** Meaningful gap exists. Consider gender-specific engagement strategies."
                        )
                    except Exception:
                        st.write("Insufficient gender data.")

    # ────────────────────────────────────────────────────────────────────────
    # TAB 4: SUBJECT WISE YoY — Grade-aware (all students, Endline only)
    # ────────────────────────────────────────────────────────────────────────
    with subj_tab:
        st.markdown("### 📚 Subject & Grade-Wise YoY Comparison (All Students, Endline Only)")
        st.caption(
            "Use the **Subject** filter to focus on one subject and see how each grade performed "
            "across both academic years. With 'All' selected, charts aggregate across all subjects."
        )
        df_sv = df_endlines.copy()

        # ── Filters row ───────────────────────────────────────────────────
        fc1, fc2 = st.columns(2)
        with fc1:
            all_subjects_sv = sorted(df_sv["Subject"].dropna().unique())
            sel_subj = st.selectbox(
                "Filter by Subject", ["All"] + list(all_subjects_sv), key="subj_sel_long",
                help="Select a single subject to see per-grade breakdowns clearly. "
                     "Some subjects only exist in one year — that's expected.",
            )
        with fc2:
            ay_opts_sv = [ay for ay in AY_ORDER if ay in df_sv["Academic Year"].values]
            sel_ay_sv  = st.selectbox(
                "Academic Year for R.I.S.E chart", ay_opts_sv,
                index=len(ay_opts_sv) - 1, key="ay_subj_long",
            )

        if sel_subj != "All":
            df_sv = df_sv[df_sv["Subject"] == sel_subj]

        # Canonical grade order (numeric sort)
        all_grades_sv = sorted(df_sv["Grade"].dropna().unique(), key=lambda x: int(x) if str(x).isdigit() else 99)

        # ── Row 1: Avg endline score by Grade (line) + RISE by Grade (bar) ─
        col_sw1, col_sw2 = st.columns(2)

        with col_sw1:
            st.markdown("#### Avg Endline Score by Grade (YoY)")
            # Each line = one Academic Year; x-axis = Grade
            grade_avg = (
                df_sv.groupby(["Grade", "Academic Year"])["Obtained Marks"]
                .mean().reset_index()
            )
            grade_avg["Grade"] = pd.Categorical(grade_avg["Grade"], categories=all_grades_sv, ordered=True)
            grade_avg = grade_avg.sort_values("Grade")

            fig_grade_line = px.line(
                grade_avg, x="Grade", y="Obtained Marks", color="Academic Year",
                markers=True, text="Obtained Marks",
                color_discrete_map=AY_COLOR_MAP,
                category_orders={"Academic Year": AY_ORDER, "Grade": all_grades_sv},
            )
            fig_grade_line.update_traces(
                textposition="top center", texttemplate="%{text:.1f}", marker=dict(size=10)
            )
            fig_grade_line.update_layout(
                xaxis_title="Grade", yaxis_title="Avg Endline Score",
                plot_bgcolor="rgba(0,0,0,0)", margin=dict(l=0, r=0, t=30),
                legend=dict(orientation="h", yanchor="top", y=-0.15, xanchor="center", x=0.5, title=""),
            )
            fig_grade_line.update_xaxes(showgrid=False, linecolor="black")
            fig_grade_line.update_yaxes(showgrid=True, gridcolor="lightgrey", zeroline=False)
            st.plotly_chart(fig_grade_line, use_container_width=True)

        with col_sw2:
            st.markdown(f"#### R.I.S.E by Grade — {sel_ay_sv}")
            if "Category" in df_sv.columns:
                df_sv_ay = df_sv[df_sv["Academic Year"] == sel_ay_sv]
                grade_cat = df_sv_ay.groupby(["Grade", "Category"]).size().reset_index(name="Count")
                grade_cat["Pct"] = grade_cat.groupby("Grade")["Count"].transform(
                    lambda x: x / x.sum() * 100
                )
                grade_cat["Grade"] = pd.Categorical(
                    grade_cat["Grade"], categories=all_grades_sv, ordered=True
                )
                grade_cat = grade_cat.sort_values("Grade")

                fig_grade_rise = px.bar(
                    grade_cat, x="Grade", y="Pct", color="Category",
                    color_discrete_map=RISE_COLORS,
                    text=grade_cat["Pct"].apply(lambda x: f"{x:.1f}%" if x > 5 else ""),
                    category_orders={"Category": RISE_ORDER, "Grade": all_grades_sv},
                )
                fig_grade_rise.update_layout(
                    barmode="stack", xaxis_title="Grade", yaxis_title="% of Students",
                    margin=dict(l=0, r=0, t=30),
                    legend=dict(orientation="h", yanchor="top", y=-0.2, xanchor="center", x=0.5, title=""),
                )
                st.plotly_chart(fig_grade_rise, use_container_width=True)

        # ── Row 2: YoY delta by Grade ─────────────────────────────────────
        st.markdown("---")
        st.markdown("#### Grade-Wise YoY Score Change (AY24-25 → AY25-26 Endline)")
        el24_grade = (
            df_sv[df_sv["Academic Year"] == "AY24-25"]
            .groupby("Grade")["Obtained Marks"].mean()
        )
        el25_grade = (
            df_sv[df_sv["Academic Year"] == "AY25-26"]
            .groupby("Grade")["Obtained Marks"].mean()
        )
        yoy_grade = (el25_grade - el24_grade).dropna().reset_index()
        yoy_grade.columns = ["Grade", "YoY Change"]
        yoy_grade["Color"] = yoy_grade["YoY Change"].apply(
            lambda x: "Improved" if x >= 0 else "Declined"
        )
        yoy_grade["Grade"] = pd.Categorical(
            yoy_grade["Grade"], categories=all_grades_sv, ordered=True
        )
        yoy_grade = yoy_grade.sort_values("Grade")

        if not yoy_grade.empty:
            fig_yoy_grade = px.bar(
                yoy_grade, x="Grade", y="YoY Change", color="Color",
                color_discrete_map={"Improved": "#00964d", "Declined": "#ed1c2d"},
                text=yoy_grade["YoY Change"].apply(lambda x: f"{x:+.2f}"),
                category_orders={"Grade": all_grades_sv},
            )
            fig_yoy_grade.add_hline(y=0, line_dash="dash", line_color="black")
            fig_yoy_grade.update_layout(
                showlegend=False, margin=dict(l=0, r=0, t=30),
                plot_bgcolor="rgba(0,0,0,0)", xaxis_title="Grade",
                yaxis_title="YoY Endline Score Change",
            )
            st.plotly_chart(fig_yoy_grade, use_container_width=True)
        else:
            st.info(
                "Both years' endline data needed for YoY grade comparison. "
                "If a Subject filter is active for a subject only present in one year, clear it."
            )

        # ── Row 3: Subject-level YoY delta (kept for completeness) ────────
        st.markdown("---")
        st.markdown("#### Subject-Wise YoY Score Change (AY24-25 → AY25-26 Endline)")
        el24s = (
            df_endlines[df_endlines["Academic Year"] == "AY24-25"]
            .groupby("Subject")["Obtained Marks"].mean()
        )
        el25s = (
            df_endlines[df_endlines["Academic Year"] == "AY25-26"]
            .groupby("Subject")["Obtained Marks"].mean()
        )
        yoy_s = (el25s - el24s).dropna().reset_index()
        yoy_s.columns = ["Subject", "YoY Change"]
        yoy_s["Color"] = yoy_s["YoY Change"].apply(lambda x: "Improved" if x >= 0 else "Declined")
        if not yoy_s.empty:
            fig_yoys = px.bar(
                yoy_s, x="Subject", y="YoY Change", color="Color",
                color_discrete_map={"Improved": "#00964d", "Declined": "#ed1c2d"},
                text=yoy_s["YoY Change"].apply(lambda x: f"{x:+.2f}"),
            )
            fig_yoys.add_hline(y=0, line_dash="dash", line_color="black")
            fig_yoys.update_layout(
                showlegend=False, margin=dict(l=0, r=0, t=30),
                plot_bgcolor="rgba(0,0,0,0)", xaxis_title="", yaxis_title="YoY Endline Score Change",
            )
            st.plotly_chart(fig_yoys, use_container_width=True)
        else:
            st.info(
                "Both years' endline data needed for YoY subject comparison. "
                "If the Subject filter is set to a subject only present in one year, clear it."
            )

    # ────────────────────────────────────────────────────────────────────────
    # TAB 5: GEOGRAPHICAL YoY (Endline only)
    # ────────────────────────────────────────────────────────────────────────
    with geo_tab:
        st.markdown("### 🗺️ Geographical YoY Comparison (Endline Only)")
        df_geo = df_endlines.copy()

        col_g1, col_g2 = st.columns(2)
        with col_g1:
            st.markdown("#### State-wise Avg Endline Score (YoY)")
            state_avg = df_geo.groupby(["State", "Academic Year"])["Obtained Marks"].mean().reset_index()
            fig_sa = px.line(
                state_avg, x="Academic Year", y="Obtained Marks", color="State",
                markers=True, text="Obtained Marks",
                category_orders={"Academic Year": AY_ORDER},
            )
            fig_sa.update_traces(textposition="top center", texttemplate="%{text:.1f}", marker=dict(size=9))
            fig_sa.update_layout(
                xaxis_title="", yaxis_title="Avg Endline Score",
                plot_bgcolor="rgba(0,0,0,0)", margin=dict(l=0, r=0, t=30),
            )
            fig_sa.update_xaxes(showgrid=False, linecolor="black")
            fig_sa.update_yaxes(showgrid=True, gridcolor="lightgrey", zeroline=False)
            st.plotly_chart(fig_sa, use_container_width=True)

        with col_g2:
            st.markdown("#### State-wise R.I.S.E Shift")
            ay_opts_geo = [ay for ay in AY_ORDER if ay in df_geo["Academic Year"].values]
            sel_ay_geo  = st.selectbox("Academic Year", ay_opts_geo, index=len(ay_opts_geo)-1, key="ay_geo_long")
            sg = df_geo[df_geo["Academic Year"] == sel_ay_geo].groupby(["State", "Category"]).size().reset_index(name="Count")
            sg["Pct"] = sg.groupby("State")["Count"].transform(lambda x: x / x.sum() * 100)
            fig_sg = px.bar(
                sg, x="State", y="Pct", color="Category",
                color_discrete_map=RISE_COLORS,
                text=sg["Pct"].apply(lambda x: f"{x:.1f}%" if x > 5 else ""),
                category_orders={"Category": RISE_ORDER},
            )
            fig_sg.update_layout(
                barmode="stack", yaxis_title="% of Students", margin=dict(l=0, r=0, t=30),
                legend=dict(orientation="h", yanchor="top", y=-0.2, xanchor="center", x=0.5, title=""),
            )
            st.plotly_chart(fig_sg, use_container_width=True)

        st.markdown("---")
        st.markdown("#### State-wise YoY Score Change")
        el24_st = df_geo[df_geo["Academic Year"] == "AY24-25"].groupby("State")["Obtained Marks"].mean()
        el25_st = df_geo[df_geo["Academic Year"] == "AY25-26"].groupby("State")["Obtained Marks"].mean()
        if not el24_st.empty and not el25_st.empty:
            yoy_st = (el25_st - el24_st).dropna().reset_index()
            yoy_st.columns = ["State", "YoY Change"]
            yoy_st["Color"] = yoy_st["YoY Change"].apply(lambda x: "Improved" if x >= 0 else "Declined")
            fig_yst = px.bar(
                yoy_st, x="State", y="YoY Change", color="Color",
                color_discrete_map={"Improved": "#00964d", "Declined": "#ed1c2d"},
                text=yoy_st["YoY Change"].apply(lambda x: f"{x:+.2f}"),
            )
            fig_yst.add_hline(y=0, line_dash="dash", line_color="black")
            fig_yst.update_layout(
                showlegend=False, margin=dict(l=0, r=0, t=30),
                plot_bgcolor="rgba(0,0,0,0)", xaxis_title="", yaxis_title="YoY Score Change",
            )
            st.plotly_chart(fig_yst, use_container_width=True)
        else:
            st.info("Both years' endline data needed for state YoY comparison.")

    # ────────────────────────────────────────────────────────────────────────
    # TAB 6: CENTRE DEEP DIVE (Endline only)
    # ────────────────────────────────────────────────────────────────────────
    with centre_tab:
        st.markdown("### 🏫 Centre Deep Dive — YoY (Endline Only)")
        df_cdd = df_endlines.copy()

        all_centres = sorted(df_cdd["Centre Name"].dropna().unique())
        sel_ctr = st.selectbox("Select Centre", all_centres, key="centre_deep_long")
        df_c = df_cdd[df_cdd["Centre Name"] == sel_ctr]

        if df_c.empty:
            st.warning("No endline data for this centre with the current filters.")
        else:
            col_c1, col_c2 = st.columns(2)
            with col_c1:
                st.markdown("#### Avg Endline Score Trajectory")
                c_avg = df_c.groupby(["Academic Year", "Subject"])["Obtained Marks"].mean().reset_index()
                fig_ca = px.line(
                    c_avg, x="Academic Year", y="Obtained Marks", color="Subject",
                    markers=True, text="Obtained Marks",
                    category_orders={"Academic Year": AY_ORDER},
                )
                fig_ca.update_traces(textposition="top center", texttemplate="%{text:.1f}", marker=dict(size=10))
                fig_ca.update_layout(
                    xaxis_title="", yaxis_title="Avg Endline Score",
                    plot_bgcolor="rgba(0,0,0,0)", margin=dict(l=0, r=0, t=30),
                )
                fig_ca.update_xaxes(showgrid=False, linecolor="black")
                fig_ca.update_yaxes(showgrid=True, gridcolor="lightgrey", zeroline=False)
                st.plotly_chart(fig_ca, use_container_width=True)

            with col_c2:
                st.markdown("#### R.I.S.E Distribution by Year")
                cc = df_c.groupby(["Academic Year", "Category"]).size().reset_index(name="Count")
                cc["Pct"] = cc.groupby("Academic Year")["Count"].transform(lambda x: x / x.sum() * 100)
                fig_cc = px.bar(
                    cc, x="Academic Year", y="Pct", color="Category",
                    color_discrete_map=RISE_COLORS,
                    text=cc["Pct"].apply(lambda x: f"{x:.1f}%" if x > 5 else ""),
                    category_orders={"Category": RISE_ORDER, "Academic Year": AY_ORDER},
                )
                fig_cc.update_layout(
                    barmode="stack", yaxis_title="% of Students", margin=dict(l=0, r=0, t=30),
                    legend=dict(orientation="h", yanchor="top", y=-0.2, xanchor="center", x=0.5, title=""),
                )
                st.plotly_chart(fig_cc, use_container_width=True)

            st.markdown("---")
            st.markdown("#### Centre vs All Centres Benchmark")
            e24c = df_c[df_c["Academic Year"] == "AY24-25"]["Obtained Marks"].mean()
            e25c = df_c[df_c["Academic Year"] == "AY25-26"]["Obtained Marks"].mean()
            e24a = df_cdd[df_cdd["Academic Year"] == "AY24-25"]["Obtained Marks"].mean()
            e25a = df_cdd[df_cdd["Academic Year"] == "AY25-26"]["Obtained Marks"].mean()
            bench = pd.DataFrame({
                "Group":         [sel_ctr, "All Centres", sel_ctr, "All Centres"],
                "Academic Year": ["AY24-25", "AY24-25", "AY25-26", "AY25-26"],
                "Avg Score":     [e24c, e24a, e25c, e25a],
            }).dropna(subset=["Avg Score"])
            if not bench.empty:
                fig_bench = px.bar(
                    bench, x="Academic Year", y="Avg Score", color="Group",
                    barmode="group",
                    text=bench["Avg Score"].apply(lambda x: f"{x:.2f}"),
                    category_orders={"Academic Year": AY_ORDER},
                )
                fig_bench.update_layout(
                    xaxis_title="", yaxis_title="Avg Endline Score", margin=dict(l=0, r=0, t=30),
                    legend=dict(orientation="h", yanchor="top", y=-0.15, xanchor="center", x=0.5, title=""),
                )
                st.plotly_chart(fig_bench, use_container_width=True)

            st.markdown("---")
            st.markdown("#### Gender Breakdown (AY 25-26 Endline)")
            gen_c = df_c[(df_c["Academic Year"] == "AY25-26") & df_c["Gender"].isin(["Boy", "Girl"])]
            if not gen_c.empty:
                gen_avg_c = gen_c.groupby("Gender")["Obtained Marks"].mean().reset_index()
                fig_gen_c = px.bar(
                    gen_avg_c, x="Gender", y="Obtained Marks", color="Gender",
                    color_discrete_map={"Boy": "#636EFA", "Girl": "#EF553B"},
                    text=gen_avg_c["Obtained Marks"].apply(lambda x: f"{x:.2f}"),
                )
                fig_gen_c.update_layout(
                    xaxis_title="", yaxis_title="Avg Score",
                    plot_bgcolor="rgba(0,0,0,0)", margin=dict(l=0, r=0, t=30), showlegend=False,
                )
                st.plotly_chart(fig_gen_c, use_container_width=True)
            else:
                st.info("No gender data for this centre in AY 25-26.")

    # ────────────────────────────────────────────────────────────────────────
    # TAB 7: STUDENT LOOKUP
    # Searches the FULL df_long (ignoring sidebar filters) so a student is
    # never hidden by an active filter.
    # ────────────────────────────────────────────────────────────────────────
    with student_tab:
        st.markdown("### 🔍 Student Lookup — Individual YoY Profile")
        st.markdown(
            "Enter a **Student ID** to view a full Year-over-Year breakdown across all "
            "assessments (Baseline & Endline for AY 24-25 and AY 25-26)."
        )
        st.caption("ℹ️ Search runs against the **full dataset** — sidebar filters are ignored here.")

        s_col, b_col = st.columns([3, 1])
        with s_col:
            raw_input = st.text_input(
                "Student ID", placeholder="e.g. 428115", key="student_search_input"
            ).strip()
        with b_col:
            st.write("")
            st.button("🔍 Search", use_container_width=True, key="student_search_btn")

        with st.expander("Or pick a student from the current filtered view", expanded=False):
            ids_in_filter = [""] + sorted(filtered_df_long["Student ID"].dropna().unique())
            picked = st.selectbox("Student IDs in current filter", ids_in_filter, key="student_pick_dd")
            if picked:
                raw_input = picked

        search_input = raw_input.strip()

        if not search_input:
            st.info("👆 Enter a Student ID above to view their full profile.")
            st.stop()

        mask       = df_long["Student ID"].astype(str).str.startswith(search_input)
        candidates = df_long[mask]["Student ID"].dropna().unique()

        if len(candidates) == 0:
            st.error(f"❌ No records found for ID starting with **{search_input}**.")
            st.stop()

        if len(candidates) > 1 and search_input not in candidates:
            exact_id = st.selectbox("Multiple matches — select exact ID:", sorted(candidates), key="exact_id_pick")
        else:
            exact_id = search_input if search_input in candidates else candidates[0]

        student_df = df_long[df_long["Student ID"] == exact_id].copy()

        if student_df.empty:
            st.error("No data found for this Student ID.")
            st.stop()

        # ── Profile header ──────────────────────────────────────────────
        st.markdown("---")
        meta = student_df.sort_values("Timepoint").iloc[0]
        grade_by_year = student_df.groupby("Academic Year")["Grade"].first().to_dict()
        gender_val = student_df[~student_df["Gender"].isin([NO_GENDER, "nan", "none", "null"])]["Gender"].mode()
        gender_display = gender_val.iloc[0] if not gender_val.empty else "Not recorded (AY 24-25 only)"

        st.markdown(
            f"""
            <div class="student-card">
                <h3 style="margin:0;color:#0094c9;">🎓 Student ID: {exact_id}</h3>
                <p style="margin:6px 0 0 0;color:#444;line-height:1.8;">
                    🏫 <b>Centre:</b> {meta.get('Centre Name','—')} &nbsp;|&nbsp;
                    🌍 <b>State:</b>  {meta.get('State','—')} &nbsp;|&nbsp;
                    🤝 <b>Donor:</b>  {meta.get('Donor','—')} &nbsp;|&nbsp;
                    🚻 <b>Gender:</b> {gender_display}<br>
                    📚 <b>Grade AY24-25:</b> {grade_by_year.get('AY24-25','—')} &nbsp;|&nbsp;
                    📚 <b>Grade AY25-26:</b> {grade_by_year.get('AY25-26','—')}
                </p>
            </div>
            """,
            unsafe_allow_html=True,
        )

        # ── KPI row ────────────────────────────────────────────────────
        el24_s = student_df[(student_df["Academic Year"] == "AY24-25") & (student_df["Period"] == "Endline")]
        el25_s = student_df[(student_df["Academic Year"] == "AY25-26") & (student_df["Period"] == "Endline")]

        k1, k2, k3, k4 = st.columns(4)
        k1.metric("Timepoints recorded", student_df["Timepoint"].nunique(),
                  help="Max possible = 4 (BL + EL × 2 years)")
        k2.metric("Subjects (AY24-25 EL)", len(el24_s))
        k3.metric("Subjects (AY25-26 EL)", len(el25_s))

        shared_subs = (
            set(el24_s["Subject"]) & set(el25_s["Subject"])
            if not (el24_s.empty or el25_s.empty) else set()
        )
        if shared_subs:
            avg24_k = el24_s[el24_s["Subject"].isin(shared_subs)]["Obtained Marks"].mean()
            avg25_k = el25_s[el25_s["Subject"].isin(shared_subs)]["Obtained Marks"].mean()
            k4.metric("Avg Score Change (shared subjects)", f"{avg25_k:.1f}",
                      delta=f"{avg25_k - avg24_k:+.1f}")
        else:
            k4.metric("YoY Status", "No shared subjects across both years")

        # ── Full assessment table ───────────────────────────────────────
        st.markdown("---")
        st.markdown("#### 📋 Full Assessment Record")
        tbl = student_df[
            ["Timepoint", "Subject", "Grade", "Obtained Marks", "Total Marks", "Pct Score", "Category"]
        ].copy().sort_values(["Timepoint", "Subject"])

        def highlight_cat(val):
            return {
                "Reviving":   "background-color:#f27c48;color:white",
                "Initiating": "background-color:#0094c9;color:white",
                "Shaping":    "background-color:#00964d;color:white",
                "Evolving":   "background-color:#ed1c2d;color:white",
            }.get(val, "")

        # BUG FIX: pandas >= 2.1 removed applymap; use map() instead
        st.dataframe(
            tbl.style.map(highlight_cat, subset=["Category"]),
            use_container_width=True, hide_index=True,
        )

        # ── Per-subject score bar charts ────────────────────────────────
        st.markdown("---")
        st.markdown("#### 📊 Score Journey by Subject")

        TP_ORDER = ["AY24-25 Baseline", "AY24-25 Endline", "AY25-26 Baseline", "AY25-26 Endline"]
        TP_COLORS = {
            "AY24-25 Baseline": "#a0aec0",
            "AY24-25 Endline":  "#636EFA",
            "AY25-26 Baseline": "#fbb6a0",
            "AY25-26 Endline":  "#00CC96",
        }
        student_df["Timepoint"] = pd.Categorical(student_df["Timepoint"], categories=TP_ORDER, ordered=True)
        subjects_found = sorted(student_df["Subject"].dropna().unique())

        if not subjects_found:
            st.info("No subject data found for this student.")
        else:
            for subj in subjects_found:
                sub_data = student_df[student_df["Subject"] == subj].sort_values("Timepoint")
                if sub_data.empty:
                    continue
                total_marks = sub_data["Total Marks"].max()
                fig_stu = go.Figure()
                for _, row in sub_data.iterrows():
                    tp = str(row["Timepoint"])
                    fig_stu.add_trace(go.Bar(
                        x=[tp], y=[row["Obtained Marks"]],
                        name=tp,
                        marker_color=TP_COLORS.get(tp, "#cccccc"),
                        text=f"<b>{row['Obtained Marks']:.0f}</b>/{total_marks:.0f}<br><i>{row['Category']}</i>",
                        textposition="outside",
                        showlegend=False,
                    ))
                fig_stu.add_hline(
                    y=total_marks, line_dash="dot", line_color="#888",
                    annotation_text=f"Max: {total_marks:.0f}", annotation_position="top right",
                )
                fig_stu.update_layout(
                    title=dict(text=f"<b>{subj}</b>", font=dict(size=15)),
                    yaxis=dict(range=[0, total_marks * 1.3], title="Marks",
                               showgrid=True, gridcolor="lightgrey"),
                    xaxis_title="",
                    plot_bgcolor="rgba(0,0,0,0)",
                    margin=dict(l=0, r=0, t=50, b=10),
                    height=300,
                )
                st.plotly_chart(fig_stu, use_container_width=True)

            # ── Category tier journey ─────────────────────────────────
            st.markdown("---")
            st.markdown("#### 🔄 R.I.S.E Tier Journey (All Subjects)")
            mig = student_df.groupby(["Timepoint", "Subject"], observed=True)["Category"].first().reset_index()
            mig = mig.sort_values("Timepoint")
            cat_num = {"Reviving": 1, "Initiating": 2, "Shaping": 3, "Evolving": 4}
            mig["Tier"] = mig["Category"].map(cat_num)
            # Cast Timepoint back to string for Plotly
            mig["Timepoint"] = mig["Timepoint"].astype(str)

            fig_mig = px.line(
                mig, x="Timepoint", y="Tier", color="Subject",
                markers=True, line_shape="linear",
                labels={"Tier": "Performance Tier", "Timepoint": ""},
                hover_data={"Category": True, "Tier": False},
                category_orders={"Timepoint": TP_ORDER},
            )
            fig_mig.update_traces(marker=dict(size=13), line=dict(width=2))
            fig_mig.update_layout(
                yaxis=dict(
                    tickvals=[1, 2, 3, 4],
                    ticktext=["🔴 Reviving", "🔵 Initiating", "🟢 Shaping", "🔴 Evolving"],
                    range=[0.5, 4.5], showgrid=True, gridcolor="#eee",
                ),
                plot_bgcolor="rgba(0,0,0,0)", margin=dict(l=0, r=0, t=30),
                xaxis=dict(showgrid=False, linecolor="black"),
                height=320,
            )
            for lvl, clr in [
                (1, "rgba(242,124,72,0.07)"),
                (2, "rgba(0,148,201,0.07)"),
                (3, "rgba(0,150,77,0.07)"),
                (4, "rgba(237,28,45,0.07)"),
            ]:
                fig_mig.add_hrect(y0=lvl-0.45, y1=lvl+0.45, fillcolor=clr, line_width=0)
            st.plotly_chart(fig_mig, use_container_width=True)

            # Improved / declined callout
            if not el24_s.empty and not el25_s.empty and shared_subs:
                improved, declined = [], []
                for sub in shared_subs:
                    sc24 = el24_s[el24_s["Subject"] == sub]["Obtained Marks"].mean()
                    sc25 = el25_s[el25_s["Subject"] == sub]["Obtained Marks"].mean()
                    (improved if sc25 >= sc24 else declined).append(f"{sub} ({sc25-sc24:+.1f})")
                if improved:
                    st.success(f"📈 **Improved (Endline YoY):** {', '.join(improved)}")
                if declined:
                    st.error(f"📉 **Declined (Endline YoY):** {', '.join(declined)}")

    st.stop()


# ==========================================
# MAIN DASHBOARD  (AY 25-26 BL vs EL)
# ==========================================
# Default fallback — used only if no file has been selected yet.
DATA_FILE_DEFAULT = "AltenUP-WB-BL-EL-Data-AY25-26.xlsx"

# Required columns every uploaded file must contain.
REQUIRED_COLS = {"State", "Centre Name", "Donor", "Subject", "Grade", "Student ID", "Obtained Marks"}
VALID_CATEGORIES = {"Reviving", "Initiating", "Shaping", "Evolving"}
VALID_PERIODS    = {"Baseline", "Endline"}

st.title("📈 Impact Analytics Dashboard")
st.markdown(
    "<p style='color:gray;font-size:1.1em;'>Comprehensive Baseline vs. Endline Performance Assessment</p>",
    unsafe_allow_html=True,
)

# ── Sidebar: nav + data management ───────────────────────────────────────────
with st.sidebar:
    try:
        st.image("evidyaloka_logo.png", width=273)
    except Exception:
        pass
    st.success(f"👤 **{st.session_state['user_first_name']}**")
    nc1, nc2 = st.columns(2)
    with nc1:
        if st.button("🏠 Home", use_container_width=True, key="nav_home_main"):
            st.session_state["current_page"] = "home"; st.rerun()
    with nc2:
        if st.button("Sign Out", use_container_width=True, key="signout_main"):
            st.session_state.update({"logged_in_email": None, "user_first_name": "User", "current_page": "home"})
            st.rerun()
    st.markdown("---")

    # ── Data Management panel ─────────────────────────────────────────────────
    st.markdown("### 📂 Data Management")

    # Auto-detect all Excel and CSV files in the same directory as app.py.
    # We look in the script's own directory so detection works identically
    # on local machines, Streamlit Cloud, and Docker deployments.
    _app_dir = os.path.dirname(os.path.abspath(__file__)) if "__file__" in dir() else os.getcwd()

    def _scan_data_files(directory: str) -> list:
        """
        Walk the repo root (non-recursive) and return a sorted list of
        all .xlsx / .xls / .csv filenames found there.
        Excludes hidden files and anything inside __pycache__ / .git dirs.
        """
        found = []
        for fname in sorted(os.listdir(directory)):
            if fname.startswith(".") or fname.startswith("~"):
                continue
            if fname.lower().endswith((".xlsx", ".xls", ".csv")):
                found.append(fname)
        return found

    available_files = _scan_data_files(_app_dir)

    # ── File selector ──────────────────────────────────────────────────────
    st.markdown("**Select dashboard data source**")

    if not available_files:
        st.error("⚠️ No Excel or CSV files found in the app directory.")
        selected_data_file = DATA_FILE_DEFAULT
    else:
        # Pre-select: honour session_state choice, else fall back to default
        default_idx = 0
        if st.session_state.get("selected_data_file") in available_files:
            default_idx = available_files.index(st.session_state["selected_data_file"])
        elif DATA_FILE_DEFAULT in available_files:
            default_idx = available_files.index(DATA_FILE_DEFAULT)

        selected_data_file = st.selectbox(
            "Choose file",
            options=available_files,
            index=default_idx,
            key="data_file_selectbox",
            help=(
                "All Excel (.xlsx/.xls) and CSV files detected in the app folder "
                "are listed here. Select one to reload the dashboard with that file."
            ),
        )

        # Detect change → clear cache and appended data so everything refreshes
        if st.session_state.get("selected_data_file") != selected_data_file:
            st.session_state["selected_data_file"] = selected_data_file
            # st.cache_data.clear() clears ALL cached functions app-wide.
            # We use this instead of load_and_prep_data.clear() because the
            # function is defined after this sidebar block, so calling it by
            # name here would raise a NameError at sidebar-render time.
            st.cache_data.clear()
            # Clear any appended data — it was validated against the old file
            for _k in ["appended_df", "append_log", "_last_upload_sig",
                       "needs_period_choice", "upload_period_choice"]:
                st.session_state.pop(_k, None)
            st.rerun()

        # Status badge
        _full_path = os.path.join(_app_dir, selected_data_file)
        if os.path.exists(_full_path):
            _size_kb = os.path.getsize(_full_path) // 1024
            st.success(f"✅ `{selected_data_file}` ({_size_kb} KB)")
        else:
            st.error(f"❌ `{selected_data_file}` not accessible.")

        # Refresh button — re-scans disk without changing selection
        if st.button("🔄 Refresh file list", use_container_width=True, key="refresh_files"):
            st.rerun()

    st.markdown("**Append additional data**")
    st.caption(
        "Upload a CSV or Excel whose columns match the selected file. "
        "Rows are validated, deduplicated, and merged instantly."
    )
    appended_file = st.file_uploader(
        "Upload CSV or Excel to append",
        type=["csv", "xlsx", "xls"],
        key="main_append_upload",
        help=(
            "Required columns: State, Centre Name, Donor, Subject, Grade, "
            "Student ID, Obtained Marks.\n\n"
            "Optional: Gender, Total Marks, Category (Rubrics), "
            "Academic Year (Baseline / Endline)."
        ),
    )
    if st.session_state.get("appended_df") is not None:
        append_rows = len(st.session_state["appended_df"])
        st.info(f"➕ **{append_rows:,} appended rows** active")
        if st.button("🗑️ Clear appended data", use_container_width=True, key="clear_append"):
            st.session_state.pop("appended_df", None)
            st.session_state.pop("append_log", None)
            st.session_state.pop("_last_upload_sig", None)
            st.session_state.pop("needs_period_choice", None)
            st.rerun()
    st.markdown("---")


# ── Core loading + cleaning functions ────────────────────────────────────────

@st.cache_data(show_spinner=False)
def load_and_prep_data(file_source):
    """Load and clean the base two-sheet Excel workbook."""
    common_cols = [
        "State", "Centre Name", "Donor", "Subject", "Grade", "Student ID",
        "Gender", "Total Marks", "Obtained Marks", "Category", "Academic Year",
    ]
    dfs = []
    try:
        xls    = pd.ExcelFile(file_source)
        sheets = xls.sheet_names
        base_sheet = "Baseline" if "Baseline" in sheets else 0
        end_sheet  = "Endline"  if "Endline"  in sheets else (1 if len(sheets) > 1 else 0)
        for sheet, label in [(base_sheet, "Baseline"), (end_sheet, "Endline")]:
            df_s = pd.read_excel(file_source, sheet_name=sheet)
            if "Rubrics" in df_s.columns:
                df_s.rename(columns={"Rubrics": "Category"}, inplace=True)
            df_s["Academic Year"] = label
            dfs.append(df_s[[c for c in common_cols if c in df_s.columns]])
    except Exception as e:
        st.error(f"Error reading Excel file: {e}")
        return pd.DataFrame()
    if not dfs:
        return pd.DataFrame()
    return _clean_main_df(pd.concat(dfs, ignore_index=True))


def _clean_main_df(df: pd.DataFrame) -> pd.DataFrame:
    """Shared cleaning pipeline used for both base and appended data."""
    for col in ["State", "Centre Name", "Donor", "Subject", "Student ID", "Gender", "Category"]:
        if col in df.columns:
            df[col] = df[col].astype(str).str.strip()
    if "Gender" in df.columns:
        df["Gender"] = df["Gender"].str.title().str.strip()
    if "Grade" in df.columns:
        df["Grade"] = df["Grade"].astype(str).str.replace(r"\.0$", "", regex=True)
    if "Category" in df.columns:
        df["Category"] = pd.Categorical(df["Category"], categories=RISE_ORDER, ordered=True)
    df["Obtained Marks"] = pd.to_numeric(df["Obtained Marks"], errors="coerce")
    return df


def _auto_assign_rise(obtained: pd.Series, total: pd.Series) -> pd.Series:
    """Assign RISE category from score percentage when Category column is absent or invalid."""
    pct = obtained / total.replace(0, pd.NA) * 100
    return pd.cut(
        pct.fillna(0),
        bins=[-1, 25, 50, 75, 101],
        labels=["Reviving", "Initiating", "Shaping", "Evolving"],
    ).astype(str)


def _parse_upload(uploaded_file) -> tuple:
    """
    Parse a user-uploaded CSV or Excel into a clean DataFrame.

    Returns (df, log_messages).

    Accepts three formats:
      1. Two-sheet Excel  — sheet 0 = Baseline, sheet 1 = Endline (same as base file)
      2. Single-sheet / CSV with Academic Year or Assessment Type column
      3. Single-sheet / CSV with no period column — rows labelled from session_state
    """
    log   = []
    fname = uploaded_file.name.lower()

    # ── Read file ─────────────────────────────────────────────────────────────
    try:
        if fname.endswith(".csv"):
            raw    = pd.read_csv(uploaded_file)
            sheets = None
        else:
            xls    = pd.ExcelFile(uploaded_file)
            sheets = xls.sheet_names
            if len(sheets) >= 2:
                dfs = []
                for sheet, label in [(sheets[0], "Baseline"), (sheets[1], "Endline")]:
                    d = pd.read_excel(uploaded_file, sheet_name=sheet)
                    if "Rubrics" in d.columns and "Category" not in d.columns:
                        d.rename(columns={"Rubrics": "Category"}, inplace=True)
                    d["Academic Year"] = label
                    dfs.append(d)
                raw = pd.concat(dfs, ignore_index=True)
                log.append(
                    f"📄 Two-sheet Excel — **{sheets[0]}** → Baseline, **{sheets[1]}** → Endline."
                )
            else:
                raw = pd.read_excel(uploaded_file, sheet_name=0)
    except Exception as e:
        return pd.DataFrame(), [f"❌ Could not read file: {e}"]

    raw.columns = raw.columns.str.strip()
    if "Rubrics" in raw.columns and "Category" not in raw.columns:
        raw.rename(columns={"Rubrics": "Category"}, inplace=True)

    # ── Required columns ──────────────────────────────────────────────────────
    missing = REQUIRED_COLS - set(raw.columns)
    if missing:
        return pd.DataFrame(), [
            f"❌ Missing required columns: **{', '.join(sorted(missing))}**.  \n"
            f"Columns in your file: `{', '.join(raw.columns)}`"
        ]

    # ── Resolve Academic Year ─────────────────────────────────────────────────
    if "Academic Year" not in raw.columns:
        if "Assessment Type" in raw.columns:
            raw["Academic Year"] = raw["Assessment Type"].astype(str).str.strip().str.title()
            log.append("ℹ️ Used **Assessment Type** column as Academic Year.")
        else:
            chosen = st.session_state.get("upload_period_choice", "Baseline")
            raw["Academic Year"] = chosen
            log.append(f"ℹ️ No period column found — all rows labelled **{chosen}**.")

    raw["Academic Year"] = raw["Academic Year"].astype(str).str.strip().str.title()
    bad_period = ~raw["Academic Year"].isin(VALID_PERIODS)
    if bad_period.any():
        bad_vals = raw.loc[bad_period, "Academic Year"].unique()[:5]
        log.append(f"⚠️ {bad_period.sum():,} rows had invalid period values (`{'`, `'.join(bad_vals)}`) — dropped.")
        raw = raw[~bad_period]

    # ── Numeric validation ────────────────────────────────────────────────────
    raw["Obtained Marks"] = pd.to_numeric(raw["Obtained Marks"], errors="coerce")
    null_marks = raw["Obtained Marks"].isna().sum()
    if null_marks:
        log.append(f"⚠️ {null_marks:,} rows had non-numeric Obtained Marks — dropped.")
        raw = raw.dropna(subset=["Obtained Marks"])

    if raw.empty:
        return pd.DataFrame(), log + ["❌ No valid rows remain after cleaning."]

    # ── Category validation / auto-assign ─────────────────────────────────────
    total_col = pd.to_numeric(raw.get("Total Marks", 10), errors="coerce").fillna(10)
    if "Category" in raw.columns:
        raw["Category"] = raw["Category"].astype(str).str.strip().str.title()
        bad_cat = ~raw["Category"].isin(VALID_CATEGORIES)
        if bad_cat.any():
            log.append(f"⚠️ {bad_cat.sum():,} rows had invalid Category — auto-assigned from score.")
            raw.loc[bad_cat, "Category"] = _auto_assign_rise(
                raw.loc[bad_cat, "Obtained Marks"], total_col[bad_cat]
            )
    else:
        raw["Category"] = _auto_assign_rise(raw["Obtained Marks"], total_col)
        log.append("ℹ️ No Category column — RISE tier auto-assigned from score %.")

    # ── Deduplication ─────────────────────────────────────────────────────────
    before = len(raw)
    raw = raw.drop_duplicates(subset=["Student ID", "Subject", "Academic Year"], keep="first")
    dupes = before - len(raw)
    if dupes:
        log.append(f"ℹ️ {dupes:,} duplicate rows removed (Student ID + Subject + Period).")

    # ── Keep dashboard columns only ───────────────────────────────────────────
    keep = ["State", "Centre Name", "Donor", "Subject", "Grade", "Student ID",
            "Gender", "Total Marks", "Obtained Marks", "Category", "Academic Year"]
    raw = raw[[c for c in keep if c in raw.columns]]
    for col, default in [("Gender", "Unknown"), ("Total Marks", 10)]:
        if col not in raw.columns:
            raw[col] = default

    log.append(f"✅ **{len(raw):,} valid rows** ready to append.")
    return _clean_main_df(raw), log


# ── Resolve which file to load ────────────────────────────────────────────────
# Use the sidebar selection; fall back to the hardcoded default.
DATA_FILE = st.session_state.get("selected_data_file", DATA_FILE_DEFAULT)
_app_dir  = os.path.dirname(os.path.abspath(__file__)) if "__file__" in dir() else os.getcwd()
_data_path = os.path.join(_app_dir, DATA_FILE)

if not os.path.exists(_data_path):
    st.error(
        f"⚠️ Selected file **`{DATA_FILE}`** not found at `{_data_path}`.  \n"
        "Use the **Data Management** panel in the sidebar to select an available file."
    )
    st.stop()

# ── Load base data ────────────────────────────────────────────────────────────
with st.spinner(f"Loading `{DATA_FILE}`…"):
    df_base = load_and_prep_data(_data_path)

if df_base.empty:
    st.error("Base data is empty. Check the file.")
    st.stop()

# ── Process new upload (only when a new file is dropped) ─────────────────────
if appended_file is not None:
    upload_sig = f"{appended_file.name}_{appended_file.size}"
    if st.session_state.get("_last_upload_sig") != upload_sig:
        st.session_state["_last_upload_sig"] = upload_sig

        # Detect whether a period selector is needed before parsing
        needs_choice = False
        try:
            fname_lower = appended_file.name.lower()
            if fname_lower.endswith(".csv"):
                probe = pd.read_csv(io.BytesIO(appended_file.read()))
                appended_file.seek(0)
                needs_choice = ("Academic Year" not in probe.columns and
                                "Assessment Type" not in probe.columns)
            else:
                probe_xls = pd.ExcelFile(io.BytesIO(appended_file.read()))
                appended_file.seek(0)
                if len(probe_xls.sheet_names) < 2:
                    probe = pd.read_excel(io.BytesIO(appended_file.read()), nrows=1)
                    appended_file.seek(0)
                    needs_choice = ("Academic Year" not in probe.columns and
                                    "Assessment Type" not in probe.columns)
        except Exception:
            pass

        st.session_state["needs_period_choice"] = needs_choice
        if needs_choice and "upload_period_choice" not in st.session_state:
            st.session_state["upload_period_choice"] = "Baseline"

        parsed_df, parse_log = _parse_upload(appended_file)
        st.session_state["append_log"]  = parse_log
        st.session_state["appended_df"] = parsed_df if not parsed_df.empty else None

# ── Period selector (shown inline when file has no period label) ──────────────
if st.session_state.get("needs_period_choice") and appended_file is not None:
    st.info(
        "Your file has no **Academic Year** or **Assessment Type** column. "
        "Select which period to assign to all uploaded rows:"
    )
    chosen_period = st.radio(
        "Assign uploaded rows to:",
        options=["Baseline", "Endline"],
        index=0 if st.session_state.get("upload_period_choice", "Baseline") == "Baseline" else 1,
        horizontal=True,
        key="upload_period_radio",
    )
    if chosen_period != st.session_state.get("upload_period_choice"):
        st.session_state["upload_period_choice"] = chosen_period
        appended_file.seek(0)
        parsed_df, parse_log = _parse_upload(appended_file)
        st.session_state["append_log"]  = parse_log
        st.session_state["appended_df"] = parsed_df if not parsed_df.empty else None
        st.rerun()

# ── Validation log ────────────────────────────────────────────────────────────
if st.session_state.get("append_log"):
    with st.expander("📋 Upload validation log", expanded=True):
        for msg in st.session_state["append_log"]:
            if msg.startswith("❌"):   st.error(msg)
            elif msg.startswith("⚠️"): st.warning(msg)
            elif msg.startswith("✅"): st.success(msg)
            else:                      st.info(msg)

# ── Merge base + appended ─────────────────────────────────────────────────────
appended_df = st.session_state.get("appended_df")
if appended_df is not None and not appended_df.empty:
    df = pd.concat([df_base, appended_df], ignore_index=True)
    df = df.drop_duplicates(subset=["Student ID", "Subject", "Academic Year"], keep="first")
    df["Category"] = pd.Categorical(df["Category"], categories=RISE_ORDER, ordered=True)
    st.success(
        f"✅ Dashboard showing base + **{len(appended_df):,} appended rows** "
        f"({len(df):,} total after deduplication). "
        f"Use **Clear appended data** in the sidebar to revert."
    )
else:
    df = df_base.copy()

if df.empty:
    st.error("No data available.")
    st.stop()

with st.sidebar:
    filtered_df, main_sel = build_filter_sidebar(df, key_prefix="main")
    selected_donors = main_sel["donor"]

if filtered_df.empty:
    st.warning("⚠️ No data for selected filters.")
    st.stop()

tab1, tab2, tab3, tab4, tab5, tab6 = st.tabs([
    "📊 Executive Summary", "📚 Subject Deep-Dive",
    "🗺️ Geographic View",  "🧑‍🎓 Student-Level Impact",
    "🚻 Gender Analysis",   "📉 RTM Analysis",
])

base_df = filtered_df[filtered_df["Academic Year"] == "Baseline"]
end_df  = filtered_df[filtered_df["Academic Year"] == "Endline"]

# ── TAB 1: EXECUTIVE SUMMARY ─────────────────────────────────
with tab1:
    st.markdown("### 🚀 High-Level Metrics")
    k1, k2, k3, k4, k5 = st.columns(5)
    matched = 0
    if not base_df.empty and not end_df.empty and "Student ID" in df.columns:
        matched = len(pd.merge(
            base_df[["Student ID","Subject"]].dropna(),
            end_df[["Student ID","Subject"]].dropna(),
            on=["Student ID","Subject"],
        ))
    avg_b = base_df["Obtained Marks"].mean() if not base_df.empty else None
    avg_e = end_df["Obtained Marks"].mean()  if not end_df.empty  else None
    sd_b  = base_df["Obtained Marks"].std()  if not base_df.empty and len(base_df)>1 else None
    sd_e  = end_df["Obtained Marks"].std()   if not end_df.empty  and len(end_df)>1  else None

    k1.metric("Matched Students", f"{matched:,}")
    if avg_b is not None and avg_e is not None:
        k2.metric("Baseline Mean", f"{avg_b:.2f}")
        k3.metric("Endline Mean",  f"{avg_e:.2f}", delta=f"{avg_e-avg_b:.2f}")
        k4.metric("Endline SD", f"{sd_e:.2f}" if sd_e else "N/A",
                  delta=f"{sd_e-sd_b:.2f}" if (sd_e and sd_b) else None, delta_color="inverse")
        be = len(base_df[base_df["Category"]=="Evolving"])/len(base_df)*100 if len(base_df) else 0
        ee = len(end_df[end_df["Category"]=="Evolving"])/len(end_df)*100   if len(end_df)  else 0
        k5.metric("Students in 'Evolving'", f"{ee:.1f}%", delta=f"{ee-be:.1f}%")
    elif avg_b:
        k2.metric("Baseline Mean",f"{avg_b:.2f}"); k3.metric("Endline Mean","N/A")
        k4.metric("Endline SD","N/A"); k5.metric("Data Status","Awaiting Endline")
    else:
        k2.metric("Baseline Mean","N/A")
        k3.metric("Endline Mean",f"{avg_e:.2f}" if avg_e else "N/A")
        k4.metric("Endline SD",  f"{sd_e:.2f}"  if sd_e  else "N/A")
        k5.metric("Data Status","Endline Only")

    st.info("**💡 SD:** A decrease means scores are clustering — the gap between high and low performers is closing.")
    st.markdown("---")
    ca, cb = st.columns(2)
    with ca:
        st.markdown("#### 📈 Score Distribution (Box Plot)")
        fig_box = px.box(filtered_df, x="Academic Year", y="Obtained Marks",
                         color="Academic Year", color_discrete_map=COLOR_MAP, points="all")
        fig_box.update_layout(showlegend=False, margin=dict(l=0,r=0,t=30))
        st.plotly_chart(fig_box, use_container_width=True)
    with cb:
        st.markdown("#### 🧬 R.I.S.E Category Shift")
        cc = filtered_df.groupby(["Academic Year","Category"]).size().reset_index(name="Count")
        cc["Pct"] = cc.groupby("Academic Year")["Count"].transform(lambda x: x/x.sum()*100)
        fig_rise = px.bar(cc, x="Category", y="Pct", color="Academic Year",
                          text=cc["Pct"].apply(lambda x: f"{x:.1f}%"),
                          color_discrete_map=COLOR_MAP, category_orders={"Category": RISE_ORDER})
        fig_rise.update_layout(barmode="group", margin=dict(l=0,r=0,t=30),
                                legend=dict(orientation="h",yanchor="top",y=-0.15,xanchor="center",x=0.5,title=""))
        st.plotly_chart(fig_rise, use_container_width=True)

# ── TAB 2: SUBJECT DEEP-DIVE ─────────────────────────────────
with tab2:
    st.markdown("### 📚 Subject & Grade Performance (R.I.S.E. Distribution)")

    def get_stacked(dfs):
        if dfs.empty or "Grade" not in dfs.columns: return pd.DataFrame()
        g = dfs.groupby(["Grade","Category"]).size().reset_index(name="Count")
        g["Pct"] = g.groupby("Grade")["Count"].transform(lambda x: x/x.sum()*100)
        return g

    bs, es = get_stacked(base_df), get_stacked(end_df)
    sc1, sc2 = st.columns(2)
    for col_w, stk, lbl in [(sc1, bs, "Baseline"), (sc2, es, "Endline")]:
        with col_w:
            st.markdown(f"#### {lbl} R.I.S.E by Grade")
            if not stk.empty:
                fig_g = px.bar(stk, x="Grade", y="Pct", color="Category",
                               color_discrete_map=RISE_COLORS,
                               text=stk["Pct"].apply(lambda x: f"{x:.1f}%" if x>5 else ""),
                               category_orders={"Category": RISE_ORDER})
                fig_g.update_layout(barmode="stack", yaxis_title="% of Students", margin=dict(l=0,r=0,t=30),
                                    legend=dict(orientation="h",yanchor="top",y=-0.15,xanchor="center",x=0.5,title=""))
                st.plotly_chart(fig_g, use_container_width=True)
            else:
                st.info(f"No {lbl} data.")

    st.markdown("---"); st.markdown("#### 🧠 Automated Insights")
    if not bs.empty and not es.empty:
        try:
            bp = bs.pivot(index="Grade",columns="Category",values="Pct").fillna(0)
            ep = es.pivot(index="Grade",columns="Category",values="Pct").fillna(0)
            for cat in RISE_ORDER:
                if cat not in bp.columns: bp[cat]=0
                if cat not in ep.columns: ep[cat]=0
            cg = bp.index.intersection(ep.index)
            if len(cg):
                dp = ep.loc[cg]-bp.loc[cg]
                beg=dp["Evolving"].idxmax(); bev=dp["Evolving"].max()
                brg=dp["Reviving"].idxmin(); brv=dp["Reviving"].min()
                if bev>0: st.success(f"📈 Grade **{beg}** had the highest shift into 'Evolving' (+{bev:.1f}pp).")
                else:     st.warning("⚠️ No grade increased its 'Evolving' share.")
                if brv<0: st.success(f"📉 Grade **{brg}** reduced 'Reviving' students most ({brv:.1f}pp).")
                else:     st.warning("⚠️ No grade reduced its 'Reviving' share.")
        except Exception:
            st.info("Not enough variance to generate insights.")
    else:
        st.info("Awaiting both Baseline and Endline data.")

# ── TAB 3: GEOGRAPHIC VIEW ────────────────────────────────────
with tab3:
    st.markdown("### 🗺️ Geographic & Centre Analysis")
    if not filtered_df.empty:
        sc = filtered_df.groupby(["State","Academic Year","Category"]).size().reset_index(name="Count")
        sc["Pct"] = sc.groupby(["State","Academic Year"])["Count"].transform(lambda x: x/x.sum()*100)
        sc["Period"] = sc["Academic Year"].map({"Baseline":"B","Endline":"E"})
        sc["St"] = sc["State"].apply(lambda s: "".join(w.upper() for w in str(s).split()) if len(str(s).split())>1 else str(s)[:3].upper())
        fig_st = px.bar(sc, x="Period", y="Pct", color="Category", facet_col="St",
                        hover_data={"State":True,"St":False,"Period":False,"Academic Year":True},
                        color_discrete_map=RISE_COLORS,
                        text=sc["Pct"].apply(lambda x: f"{x:.1f}%" if x>5 else ""),
                        category_orders={"Category":RISE_ORDER,"Period":["B","E"]})
        fig_st.update_layout(barmode="stack",yaxis_title="% of Students",margin=dict(l=0,r=0,t=40),
                             legend=dict(orientation="h",yanchor="top",y=-0.2,xanchor="center",x=0.5,title=""))
        fig_st.update_xaxes(title_text="")
        fig_st.for_each_annotation(lambda a: a.update(text=a.text.split("=")[-1]))
        st.plotly_chart(fig_st, use_container_width=True)

    st.markdown("---"); st.markdown("#### Top 10 Centres (by % Evolving)")
    if not filtered_df.empty:
        ccat = filtered_df.groupby(["Centre Name","Category"]).size().reset_index(name="Count")
        ccat["Pct"] = ccat.groupby("Centre Name")["Count"].transform(lambda x: x/x.sum()*100)
        cp = ccat.pivot(index="Centre Name",columns="Category",values="Pct").fillna(0)
        for cat in RISE_ORDER:
            if cat not in cp.columns: cp[cat]=0
        cp = cp.sort_values(by=["Evolving","Shaping","Initiating","Reviving"],ascending=False).head(10).iloc[::-1]
        tc = cp.reset_index().melt(id_vars="Centre Name",value_vars=RISE_ORDER,var_name="Category",value_name="Pct")
        fig_tc = px.bar(tc,x="Pct",y="Centre Name",color="Category",orientation="h",
                        color_discrete_map=RISE_COLORS,
                        text=tc["Pct"].apply(lambda x: f"{x:.1f}%" if x>5 else ""),
                        category_orders={"Category":RISE_ORDER})
        fig_tc.update_layout(barmode="stack",xaxis_title="% of Students",yaxis_title="",margin=dict(l=0,r=0,t=30),
                             legend=dict(orientation="h",yanchor="top",y=-0.15,xanchor="center",x=0.5,title=""))
        st.plotly_chart(fig_tc, use_container_width=True)

# ── TAB 4: STUDENT-LEVEL IMPACT ──────────────────────────────
with tab4:
    st.markdown("### 🧑‍🎓 Student-Level Impact (Matched Cohort)")
    if not base_df.empty and not end_df.empty and "Student ID" in df.columns:
        bc = base_df[["Student ID","Subject","Obtained Marks","Category"]].dropna(subset=["Student ID"]).drop_duplicates(subset=["Student ID","Subject"])
        ec = end_df[["Student ID","Subject","Obtained Marks","Category"]].dropna(subset=["Student ID"]).drop_duplicates(subset=["Student ID","Subject"])
        paired = pd.merge(bc, ec, on=["Student ID","Subject"], suffixes=("_BL","_EL"))
        if not paired.empty:
            paired["Delta"] = paired["Obtained Marks_EL"]-paired["Obtained Marks_BL"]
            tot=len(paired); mc=paired["Delta"].mean()
            pp=len(paired[paired["Delta"]>0])/tot*100
            np_=len(paired[paired["Delta"]==0])/tot*100
            mn=len(paired[paired["Delta"]<0])/tot*100
            m1,m2,m3,m4,m5=st.columns(5)
            m1.metric("Matched",f"{tot:,}"); m2.metric("Avg Change",f"{mc:+.2f}")
            m3.metric("Improved",f"{pp:.1f}%"); m4.metric("Unchanged",f"{np_:.1f}%"); m5.metric("Declined",f"{mn:.1f}%")
            st.markdown("---"); st.markdown("#### 🔄 Category Transition Matrix")
            st.caption("Rows = Baseline. Columns = Endline. Green = upward, grey = no change, red = downward.")
            tm = pd.crosstab(paired["Category_BL"],paired["Category_EL"],normalize="index")*100
            tm = tm.reindex(index=RISE_ORDER,columns=RISE_ORDER,fill_value=0)
            dm = pd.DataFrame(
                [[0 if i==j else (1 if j>i else -1) for j in range(4)] for i in range(4)],
                index=RISE_ORDER, columns=RISE_ORDER, dtype=float,
            )
            fig_heat = px.imshow(dm, x=tm.columns, y=tm.index,
                                 color_continuous_scale=["#FF7F7F","#F2F4F7","#82E0AA"],
                                 labels=dict(x="Endline Category",y="Baseline Category"))
            fig_heat.update_traces(text=tm.map(lambda x: f"{x:.1f}%"), texttemplate="%{text}",
                                   hovertemplate="BL: %{y}<br>EL: %{x}<br>%{text}<extra></extra>")
            fig_heat.update_coloraxes(showscale=False)
            fig_heat.update_layout(margin=dict(l=0,r=0,t=30,b=0),height=500)
            _, hc, _ = st.columns(3)
            with hc: st.plotly_chart(fig_heat, use_container_width=True)
        else:
            st.warning("No matched Student ID + Subject pairs between BL and EL.")
    else:
        st.info("Both Baseline and Endline with Student ID needed.")

# ── TAB 5: GENDER ANALYSIS ───────────────────────────────────
with tab5:
    st.markdown("### 🚻 Gender-Wise Performance")
    if "Gender" in filtered_df.columns:
        gdf = filtered_df[~filtered_df["Gender"].astype(str).str.lower().isin(["nan","none","null",""])].copy()
        if not gdf.empty:
            gb = gdf[gdf["Academic Year"]=="Baseline"]
            ge = gdf[gdf["Academic Year"]=="Endline"]
            gens = sorted(gdf["Gender"].dropna().unique())
            gcols = st.columns(max(len(gens),2))
            for i,g in enumerate(gens):
                bm = gb[gb["Gender"]==g]["Obtained Marks"].mean() if not gb.empty else None
                em = ge[ge["Gender"]==g]["Obtained Marks"].mean() if not ge.empty else None
                with gcols[i]:
                    if bm and em: st.metric(f"{g} Endline",f"{em:.2f}",delta=f"{em-bm:.2f}")
                    elif em:      st.metric(f"{g} Endline",f"{em:.2f}")
                    elif bm:      st.metric(f"{g} Baseline",f"{bm:.2f}")
            st.markdown("---")
            gc1,gc2=st.columns(2)
            with gc1:
                ag = gdf.groupby(["Gender","Academic Year"])["Obtained Marks"].mean().reset_index()
                fig_ga = px.bar(ag,x="Gender",y="Obtained Marks",color="Academic Year",
                                barmode="group",color_discrete_map=COLOR_MAP,text_auto=".2f")
                fig_ga.update_layout(yaxis_title="Avg Marks",margin=dict(l=0,r=0,t=30),
                                     legend=dict(orientation="h",yanchor="top",y=-0.15,xanchor="center",x=0.5,title=""))
                st.plotly_chart(fig_ga,use_container_width=True)
            with gc2:
                gc_d = gdf.groupby(["Gender","Academic Year","Category"]).size().reset_index(name="Count")
                gc_d["Pct"] = gc_d.groupby(["Gender","Academic Year"])["Count"].transform(lambda x: x/x.sum()*100)
                fig_gc = px.bar(gc_d,x="Academic Year",y="Pct",color="Category",facet_col="Gender",
                                color_discrete_map=RISE_COLORS,
                                text=gc_d["Pct"].apply(lambda x: f"{x:.1f}%" if x>5 else ""),
                                category_orders={"Category":RISE_ORDER,"Academic Year":["Baseline","Endline"]})
                fig_gc.update_layout(barmode="stack",yaxis_title="% of Students",margin=dict(l=0,r=0,t=40),
                                     legend=dict(orientation="h",yanchor="top",y=-0.2,xanchor="center",x=0.5,title=""))
                fig_gc.update_xaxes(title_text="")
                fig_gc.for_each_annotation(lambda a: a.update(text=a.text.split("=")[-1]))
                st.plotly_chart(fig_gc,use_container_width=True)
        else:
            st.info("No valid gender data in current selection.")
    else:
        st.warning("'Gender' column missing from dataset.")

# ── TAB 6: RTM ANALYSIS ──────────────────────────────────────
with tab6:
    st.markdown("### 📉 Regression to the Mean (RTM) Analysis")
    if not base_df.empty and not end_df.empty and "Student ID" in df.columns:
        br = base_df[["Student ID","Subject","Obtained Marks"]].dropna(subset=["Student ID","Obtained Marks"]).drop_duplicates(subset=["Student ID","Subject"])
        er = end_df[["Student ID","Subject","Obtained Marks"]].dropna(subset=["Student ID","Obtained Marks"]).drop_duplicates(subset=["Student ID","Subject"])
        rtm = pd.merge(br,er,on=["Student ID","Subject"],suffixes=("_BL","_EL"))
        if not rtm.empty:
            norm = st.checkbox("⚙️ Normalise scores (Z-scores)",value=False)
            if norm:
                for col in ["Obtained Marks_BL","Obtained Marks_EL"]:
                    rtm[col]=(rtm[col]-rtm[col].mean())/rtm[col].std()
            rtm["Delta"]=rtm["Obtained Marks_EL"]-rtm["Obtained Marks_BL"]
            corr=rtm["Obtained Marks_BL"].corr(rtm["Delta"])
            var=rtm["Obtained Marks_BL"].var()
            cov=rtm["Obtained Marks_BL"].cov(rtm["Delta"])
            slope=cov/var if var and not pd.isna(var) else 0.0
            intercept=rtm["Delta"].mean()-(slope*rtm["Obtained Marks_BL"].mean()) if not pd.isna(slope) else 0.0
            tot=len(rtm)
            imp=len(rtm[rtm["Delta"]>0])/tot*100; dec=len(rtm[rtm["Delta"]<0])/tot*100
            tag=("Strong RTM" if slope<=-0.3 else "Moderate RTM" if slope<=-0.1 else "Minimal RTM" if slope<0 else "No RTM")
            r1,r2,r3,r4=st.columns(4)
            r1.metric("Correlation (r)",f"{corr:.3f}" if not pd.isna(corr) else "N/A")
            r2.metric("Slope",f"{slope:.3f}"); r3.metric("Impr/Decl",f"{imp:.1f}%/{dec:.1f}%"); r4.metric("Interpretation",tag)
            (st.warning if slope<=-0.1 else st.success)(
                "💡 Part of improvement may be statistical RTM." if slope<=-0.1
                else "💡 Growth is more likely attributable to actual intervention impact."
            )
            st.markdown("---")
            fig_rtm = px.scatter(rtm,x="Obtained Marks_BL",y="Delta",trendline="ols",
                                 trendline_color_override="red",opacity=0.6,
                                 color_discrete_sequence=["#636EFA"],
                                 labels={"Obtained Marks_BL":"Baseline Score","Delta":"Score Delta"})
            fig_rtm.add_hline(y=0,line_dash="dash",line_color="black",
                              annotation_text="No Change",annotation_position="bottom right")
            fig_rtm.update_layout(margin=dict(l=0,r=0,t=30))
            st.plotly_chart(fig_rtm,use_container_width=True)
            try:
                rtm["Quintile"]=pd.qcut(rtm["Obtained Marks_BL"],q=5,duplicates="drop")
                bstats=rtm.groupby("Quintile",observed=False).agg(
                    Avg_BL=("Obtained Marks_BL","mean"),
                    Avg_Delta=("Delta","mean"),
                    N=("Student ID","count"),
                ).reset_index()
                bstats["Q_str"]=bstats["Quintile"].astype(str)
                bstats=bstats.sort_values("Avg_BL")
                fig_bin=px.bar(bstats,x="Q_str",y="Avg_Delta",
                               text=bstats["Avg_Delta"].apply(lambda x: f"{x:+.2f}"),
                               color="Avg_Delta",color_continuous_scale=px.colors.diverging.RdYlGn,
                               color_continuous_midpoint=0,
                               labels={"Q_str":"Baseline Quintile","Avg_Delta":"Avg Delta"})
                fig_bin.add_hline(y=0,line_dash="solid",line_color="black",line_width=1)
                fig_bin.update_traces(textposition="outside")
                fig_bin.update_layout(margin=dict(l=0,r=0,t=30,b=40),coloraxis_showscale=False)
                st.plotly_chart(fig_bin,use_container_width=True)
            except ValueError:
                st.info("Not enough variance for quintile bins.")
            r2_val=corr**2 if not pd.isna(corr) else 0
            s1,s2,s3=st.columns(3)
            s1.metric("r",f"{corr:.3f}" if not pd.isna(corr) else "N/A")
            s2.metric("Slope",f"{slope:.3f}"); s3.metric("R²",f"{r2_val:.3f}")
            eq=f"**Delta = {intercept:.2f} + ({slope:.2f} × Baseline)**"
            if slope<-0.3:   st.success(f"✔️ Strong RTM confirmed. {eq}")
            elif slope<-0.1: st.info(f"ℹ️ Moderate RTM. {eq}")
            elif slope<0:    st.warning(f"⚠️ Weak RTM. {eq}")
            else:            st.error(f"❌ No RTM (slope positive). {eq}")
        else:
            st.warning("No matched Student ID pairs for RTM analysis.")
    else:
        st.info("Both Baseline and Endline with Student ID column needed.")

# ── PPTX REPORT ──────────────────────────────────────────────
with st.sidebar:
    st.markdown("---"); st.markdown("### 📄 DRM Compliance Report")
    if selected_donors != "All":
        report_name = f"AY25-26_Impact_Report_{selected_donors.replace(' ','_')}.pptx"
        if st.session_state.get("ready_ppt_donor") != selected_donors:
            st.session_state.pop("ready_ppt", None); st.session_state.pop("ready_ppt_donor", None)

        if st.button(f"⚙️ Prepare PPTX for {selected_donors}", use_container_width=True):
            with st.spinner("Generating presentation…"):
                try:
                    from pptx import Presentation
                    from pptx.util import Inches
                    prs = Presentation()
                    chart_figs = {}

                    def fig_to_png(fig):
                        fig.update_layout(plot_bgcolor="white", paper_bgcolor="white")
                        buf = io.BytesIO()
                        fig.write_image(buf, format="png", engine="kaleido", width=1000, height=550)
                        buf.seek(0); return buf

                    def add_slide(fig, title):
                        sl = prs.slides.add_slide(prs.slide_layouts[5])
                        sl.shapes.title.text = title
                        sl.shapes.add_picture(fig_to_png(fig), Inches(0.5), Inches(1.5), width=Inches(9))

                    sl1 = prs.slides.add_slide(prs.slide_layouts[0])
                    sl1.shapes.title.text = "AY 25-26 Impact Report"
                    sl1.placeholders[1].text = f"Donor: {selected_donors}"

                    sl2 = prs.slides.add_slide(prs.slide_layouts[1])
                    sl2.shapes.title.text = "Executive Summary"
                    tf = sl2.placeholders[1].text_frame; tf.word_wrap = True
                    tf.text = f"Centres: {filtered_df['Centre Name'].nunique()}"
                    tf.add_paragraph().text = f"States: {', '.join(str(s) for s in sorted(filtered_df['State'].dropna().unique()))}"
                    tf.add_paragraph().text = f"Subjects: {', '.join(sorted(filtered_df['Subject'].dropna().unique()))}"

                    _cc = filtered_df.groupby(["Academic Year","Category"]).size().reset_index(name="Count")
                    _cc["Pct"] = _cc.groupby("Academic Year")["Count"].transform(lambda x: x/x.sum()*100)
                    chart_figs["rise"] = px.bar(_cc, x="Category", y="Pct", color="Academic Year",
                                                color_discrete_map=COLOR_MAP, category_orders={"Category":RISE_ORDER})
                    chart_figs["rise"].update_layout(barmode="group")
                    chart_figs["box"] = px.box(filtered_df, x="Academic Year", y="Obtained Marks",
                                               color="Academic Year", color_discrete_map=COLOR_MAP)
                    _gdf2 = filtered_df[~filtered_df["Gender"].astype(str).str.lower().isin(["nan","none","null",""])]
                    if not _gdf2.empty:
                        _ag2 = _gdf2.groupby(["Gender","Academic Year"])["Obtained Marks"].mean().reset_index()
                        chart_figs["gender"] = px.bar(_ag2, x="Gender", y="Obtained Marks",
                                                      color="Academic Year", barmode="group",
                                                      color_discrete_map=COLOR_MAP)

                    add_slide(chart_figs["rise"], "Overall R.I.S.E Shift")
                    add_slide(chart_figs["box"],  "Score Distribution")
                    if "gender" in chart_figs:
                        add_slide(chart_figs["gender"], "Gender Performance")

                    for subj in sorted(filtered_df["Subject"].dropna().unique()):
                        for period_df2, lbl2 in [(base_df,"Baseline"),(end_df,"Endline")]:
                            sd = period_df2[period_df2["Subject"]==subj] if "Subject" in period_df2.columns else pd.DataFrame()
                            if not sd.empty and "Grade" in sd.columns:
                                grp = sd.groupby(["Grade","Category"]).size().reset_index(name="Count")
                                grp["Pct"] = grp.groupby("Grade")["Count"].transform(lambda x: x/x.sum()*100)
                                _f = px.bar(grp,x="Grade",y="Pct",color="Category",
                                            color_discrete_map=RISE_COLORS,barmode="stack",
                                            category_orders={"Category":RISE_ORDER})
                                add_slide(_f, f"{lbl2} R.I.S.E — {subj}")

                    buf = io.BytesIO(); prs.save(buf); buf.seek(0)
                    st.session_state["ready_ppt"] = buf.getvalue()
                    st.session_state["ready_ppt_donor"] = selected_donors
                    st.success("Report ready!")
                except ImportError:
                    st.error("Install python-pptx and kaleido first.")
                except Exception as e:
                    st.error(f"Error generating report: {e}")

        if "ready_ppt" in st.session_state and st.session_state.get("ready_ppt_donor") == selected_donors:
            st.download_button(
                "⬇️ Download Presentation",
                data=st.session_state["ready_ppt"],
                file_name=report_name,
                mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
                use_container_width=True,
            )
    else:
        st.info("💡 Select a specific Donor to enable the DRM Report generator.")
